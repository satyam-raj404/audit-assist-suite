import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
import tkinter as tk
from tkinter import filedialog, messagebox
from PIL import Image, ImageTk
import os
from pathlib import Path

# ------------------ Utility Functions ------------------
def delete_rows_with_only_dashes(table):
    """
    Deletes rows from a PowerPoint table where all cells are empty or contain only a dash ("-").
    
    Parameters:
        table (pptx.table.Table): The table object to process.
    """
    rows_to_delete = []

    for row_idx, row in enumerate(table.rows):
        # Check each cell starting from index 1 manually
        all_dashes = True
        for cell_idx in range(0, len(row.cells)):
            if row.cells[cell_idx].text.strip() not in ["", "-"]:
                all_dashes = False
                break
        if all_dashes:
            rows_to_delete.append(row_idx)

    # Delete rows in reverse order to avoid index shifting
    for idx in reversed(rows_to_delete):
        table._tbl.remove(table._tbl.tr_lst[idx])

def set_column_width(table, col_index, width_in_inches):
    """
    Sets the width of a specific column in a PowerPoint table.

    Parameters:
        table (pptx.table.Table): The table object.
        col_index (int): Index of the column to adjust.
        width_in_inches (float): Desired width in inches.
    """
    width = Inches(width_in_inches)
    table.columns[col_index].width = width


def delete_empty_rows(table):
    # Remove rows where first cell is still "-"
    tbl = table._tbl
    # rows = list(table.rows)
    i = 0
    for row in table.rows:  # Exclude total row
        if row.cells[0].text == "-":
            tr = tbl.tr_lst[i]
            tbl.remove(tr)
        i += 1

def delete_last_column_if_empty(table):
    """
    Deletes the last column of a PowerPoint table if all cells (except the first row)
    are empty or contain only a dash ("-").
    
    Parameters:
        table (pptx.table.Table): The table object to process.
    """
    last_col_idx = len(table.columns) - 1

    set_column_width(table, last_col_idx, 0)

    # Check if all cells in the last column (excluding the header row) are empty or "-"
    all_empty = True
    for i in range(1, len(table.rows)):  # Skip the header row
        cell_text = table.cell(i, last_col_idx).text.strip()
        if cell_text not in ["", "-"]:
            all_empty = False
            break

    # Delete the last column if condition is met
    if all_empty:
        for row in table._tbl.tr_lst:
            row.remove(row.tc_lst[last_col_idx])

def delete_last_row_if_empty(table):
    """
    Deletes the last row of a PowerPoint table if all cells except the first one
    are empty or contain only a dash ("-").
    
    Parameters:
        table (pptx.table.Table): The table object to process.
    """
    last_row_idx = len(table.rows) - 1
    last_row = table.rows[last_row_idx]

    # Check if all cells except the first one are empty or contain "-"
    all_empty = True
    for i in range(1, len(last_row.cells)):
        cell_text = last_row.cells[i].text.strip()
        if cell_text not in ["", "-"]:
            all_empty = False
            break

    # Delete the last row if condition is met
    if all_empty:
        tbl = table._tbl
        tr = tbl.tr_lst[last_row_idx]
        tbl.remove(tr)                    

def duplicate_slide(pres, slide):
    """Duplicate a slide from a presentation."""
    slide_layout = pres.slide_layouts[6]  # Use blank layout
    new_slide = pres.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        el = shape.element
        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
    return new_slide

def fill_cell(cell, text, font_size=12, bold=False, align=PP_ALIGN.CENTER, color=RGBColor(255, 255, 255)):
    """Helper to fill a table cell with formatted text."""
    cell.text = str(text)
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.alignment = align

def fill_cell_b(cell, text, font_size=12, bold=False, align=PP_ALIGN.CENTER, color=RGBColor(0, 0, 0)):
    """Helper to fill a table cell with formatted text."""
    cell.text = str(text)
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.alignment = align

def fill_cell_b_lalign(cell, text, font_size=12, bold=False, align=PP_ALIGN.LEFT, color=RGBColor(0, 0, 0)):
    """Helper to fill a table cell with formatted text."""
    cell.text = str(text)
    cell.margin_left = Pt(5)
    paragraph = cell.text_frame.paragraphs[0]
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    paragraph.alignment = align

def center_table_on_slide(table_shape, prs):
    """
    Centers a table shape on the slide using presentation dimensions.

    Parameters:
        table_shape (pptx.shapes.graphfrm.GraphicFrame): The shape containing the table.
        prs (pptx.Presentation): The presentation object.
    """
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    table_width = table_shape.width
    table_height = table_shape.height

    table_shape.left = int((slide_width - table_width) / 2)
    # table_shape.top = int((slide_height - table_height) / 2)

def center_table_between_shapes(table_shape, top_shape, bottom_shape):
    """
    Centers the table vertically between two given shapes.

    Parameters:
    - table_shape: The shape object representing the table.
    - top_shape: The shape object above the table.
    - bottom_shape: The shape object below the table.
    """
    # Calculate the vertical center between the bottom of the top shape and the top of the bottom shape
    top_bottom = top_shape.top + top_shape.height
    bottom_top = bottom_shape.top
    center_y = (top_bottom + bottom_top) / 2

    # Adjust the table's top position to center it vertically
    table_shape.top = int(center_y - (table_shape.height / 2))

def append_hyperlink_to_notes_in_slide(slide, link_text, url):
    """
    Appends a clickable hyperlink below the text 'Notes' in a slide's text box.

    Parameters:
    - slide: pptx.slide.Slide object
    - link_text: str, the display text for the hyperlink
    - url: str, the URL to link to
    """
    for shape in slide.shapes:
        if shape.has_text_frame and "Notes:" in shape.text:
            # Append a new paragraph with hyperlink
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = link_text
            run.font.size = Pt(12)

            # Create hyperlink relationship
            r_id = slide.part.relate_to(
                url,
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
                is_external=True
            )

            # Apply hyperlink to the run with correct namespace
            rPr = run._r.get_or_add_rPr()
            hlinkClick = parse_xml(
                f'<a:hlinkClick xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:id="{r_id}"/>'
            )
            rPr.append(hlinkClick)

# ------------------ Main Logic ------------------

def Concurrent_Dashboard_Main(tracker_path, template_path, output_folder_path):
    # Load Excel data
    df_raw = pd.read_excel(tracker_path,sheet_name='Issues', header=None, engine='openpyxl')

    # Apply forward fill only to the first row
    df_raw.iloc[0] = df_raw.iloc[0].ffill()

    # Extract the first two rows which may contain merged header information
    header_row_1 = df_raw.iloc[0]
    header_row_2 = df_raw.iloc[1]

    # Combine the two header rows to create complete column names
    combined_headers = []
    for h1, h2 in zip(header_row_1, header_row_2):
        h1_str = str(h1) if pd.notna(h1) else ''
        h2_str = str(h2) if pd.notna(h2) else ''
        if h1_str and h2_str:
            combined_headers.append(f"{h1_str} - {h2_str}")
        elif h1_str:
            combined_headers.append(h1_str)
        elif h2_str:
            combined_headers.append(h2_str)
        else:
            combined_headers.append("Unnamed")

    # Apply the new headers and drop the first two rows
    df = df_raw.iloc[2:].reset_index(drop=True)
    df.columns = combined_headers
    
    reports_list = df['Name of the report'].dropna().unique().tolist()

    # Load template and create new presentation
    template_ppt = Presentation(template_path)
    new_ppt = Presentation()
    new_ppt.slide_width = template_ppt.slide_width
    new_ppt.slide_height = template_ppt.slide_height

    # Duplicate the first slide from template
    default_template = template_ppt.slides[0]
    new_slide = duplicate_slide(new_ppt, default_template)

    # Find the first table in the slide
    table = next((shape.table for shape in new_slide.shapes if shape.has_table), None)
    if not table:
        raise ValueError("No table found in the slide.")

    # Group by 'Name of the report' and get unique 'Month-Year' values
    report_months = df.groupby("Name of the report")["Month-Year"].unique()
    # Count unique months for each report type
    report_month_counts = {report: len(months) for report, months in report_months.items()}

    # Fill data for each report
    for report in reports_list:
        for row in table.rows:
            if row.cells[1].text == "-":
                fill_cell_b(row.cells[1], report_month_counts[report], font_size=14,)

            # Define column indices
            columns = {
                'Risk Category Tagging - H': df.columns.get_loc("Risk Category Tagging - H"),
                'Risk Category Tagging - M': df.columns.get_loc("Risk Category Tagging - M"),
                'Risk Category Tagging - L': df.columns.get_loc("Risk Category Tagging - L"),
                'Status': df.columns.get_loc("Status"),
                'Draft/Updated Draf/Final': df.columns.get_loc("Draft/Updated Draf/Final")
            }

            total = 0
            for i, col in enumerate(['Risk Category Tagging - H', 'Risk Category Tagging - H', 'Risk Category Tagging - M', 'Risk Category Tagging - M', 'Risk Category Tagging - L', 'Risk Category Tagging - L']):
                status = "Open" if i % 2 == 0 else "Closed"
                filtered = df[
                    (df['Name of the report'] == report) &
                    (df[df.columns[columns[col]]] == "✓") &
                    (df[df.columns[columns['Status']]] == status) &
                    (df[df.columns[columns['Draft/Updated Draf/Final']]] == "Final")
                ]
                if row.cells[i + 2].text == "-":
                    fill_cell_b(row.cells[i + 2], len(filtered), font_size=14)
                total += len(filtered)

            # Fill total column
            if row.cells[8].text == "-":
                fill_cell_b(row.cells[8], total, font_size=12)

            # Fill report name
            if row.cells[0].text == "-":
                fill_cell_b_lalign(row.cells[0], report, font_size=12, bold=True, align=PP_ALIGN.LEFT)
                break  # Fill only one row per slide

    # Fill total row (assumed to be row index 11)
    total_row = table.rows[11]

    # Calculate sum of column index 1 from rows 0 to 10
    column_index = 1
    column_sum = 0
    for i in range(11):
        cell_text = table.rows[i].cells[column_index].text
        try:
            column_sum += int(cell_text)
        except ValueError:
            continue  # Skip non-integer or empty cells

    # Fill the first cell in the total row with the calculated sum
    fill_cell(total_row.cells[column_index], column_sum, font_size=12, bold=True) 

    f_total = 0
    for i, col in enumerate(['Risk Category Tagging - H', 'Risk Category Tagging - H', 'Risk Category Tagging - M', 'Risk Category Tagging - M', 'Risk Category Tagging - L', 'Risk Category Tagging - L']):
        status = "Open" if i % 2 == 0 else "Closed"
        total_yes_count = sum(
            len(df[
                (df['Name of the report'] == report) &
                (df[df.columns[columns[col]]] == "✓") &
                (df[df.columns[columns['Status']]] == status) &
                (df[df.columns[columns['Draft/Updated Draf/Final']]] == "Final")
            ])
            for report in reports_list
        )
        f_total += total_yes_count
        fill_cell(total_row.cells[i + 2], total_yes_count, bold=True)

    # Fill final total
    fill_cell(total_row.cells[8], f_total, bold=True)
    
    # To delete unnecessary rows and columns:
    delete_rows_with_only_dashes(table)
    delete_last_row_if_empty(table)
    delete_last_column_if_empty(table)

    # Center aligning the table
    for shape in new_slide.shapes:
        if shape.has_table:
            center_table_on_slide(shape, new_ppt)  # Pass the shape, not shape.table
            break
    
    # Vertically aligning the table
    
    # Identify the table shape
    table_shape = None
    for shape in new_slide.shapes:
        if shape.shape_type == 19: # Table shape type
            table_shape = shape
            break

    # If table is found, find shapes above and below it
    top_shape = None
    bottom_shape = None

    if table_shape:
        table_top = table_shape.top
        table_bottom = table_shape.top + table_shape.height

        for shape in new_slide.shapes:
            if shape == table_shape:
                 continue
            shape_top = shape.top
            shape_bottom = shape.top + shape.height

            # Check if shape is above the table
            if shape_bottom <= table_top:
                if top_shape is None or shape_bottom > (top_shape.top + top_shape.height):
                    top_shape = shape

            # Check if shape is below the table
            if shape_top >= table_bottom:
                if bottom_shape is None or shape_top < bottom_shape.top:
                    bottom_shape = shape
    
    center_table_between_shapes(table_shape, top_shape, bottom_shape)

    append_hyperlink_to_notes_in_slide(new_slide,"PowerBI Dashboard", "https://app.powerbi.com/view?r=eyJrIjoiYjNkNGUxYjktZTRlNi00OTYyLWFmZGQtYjBjYmY4NTA2Yzc2IiwidCI6ImNlMjVmZjc0LWRhMDktNGRjOC1hNWZkLTVhOTI3MzkzNWRmNCIsImMiOjEwfQ%3D%3D")

    # Save the final presentation
    
    # Define the output folder inside the parent directory
    output_folder_path = Path(output_folder_path)
    output_folder = output_folder_path / "Concurrent Report Dashboard"
    # Create the folder if it doesn't exist
    output_folder.mkdir(parents=True, exist_ok=True)

    OP_DIR = output_folder
    if not os.path.exists(OP_DIR):
            os.makedirs(OP_DIR)
    output_path = r"%s/Concurrent Audit Dashboard.pptx" % (OP_DIR)
    new_ppt.save(output_path)
    return output_path
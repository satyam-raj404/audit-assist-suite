import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
import copy
from copy import deepcopy
from pptx.oxml.ns import qn
from openpyxl import load_workbook
import tkinter as tk
from tkinter import filedialog, messagebox
import os
from pathlib import Path
import shutil
import subprocess
from PIL import Image, ImageTk
import time

# Load Excel data
df = pd.read_excel("DIB - ICOFR 2024 - Detailed Gap Report- Automation Requirement.xlsx", engine="openpyxl")

# Load presentation and locate chart
prs = Presentation(r"C:\Users\aryansharma8\OneDrive - KPMG\Documents\POC\FS Automation\ICOFR\DIB - ICOFR 2024 - Detailed Gap Report - Detailed Observations Template.pptx")

new_slide = prs.slides[0]

shapes_with_tables = [shape for shape in new_slide.shapes if shape.has_table]

table_idx = 0
for shape in shapes_with_tables:
    if shape.has_table:
        table = shape.table
        row_idx = 0
        for row in table.rows:
            cell_idx = 0
            for cell in row.cells:
                text_cell = cell.text
                cell.text = text_cell + ", Table index: " + str(table_idx) + ", Row index: " + str(row_idx) + ", Cell index: " + str(cell_idx)
                cell_idx += 1
                    # Apply formatting to all paragraphs and runs in the cell
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Univers for KPMG"
                        run.font.size = Pt(11)
                        run.font.color.rgb = RGBColor(0,0,0) #002060
            row_idx += 1
        table_idx += 1 

# Save updated presentation
prs.save("Updated_Dashboard_PreservedSize.pptx")

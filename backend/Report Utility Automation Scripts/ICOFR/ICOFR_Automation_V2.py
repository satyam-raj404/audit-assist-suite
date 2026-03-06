import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
import copy
from copy import deepcopy
from pptx.oxml.ns import qn
from openpyxl import load_workbook
import tkinter as tk
from tkinter import filedialog, messagebox
import os
from pathlib import Path
import shutil
import subprocess
from PIL import Image, ImageTk
import time

# Function to duplicate a slide by copying its XML element tree
def duplicate_slide(ppt, slide):
    # Create a new slide using the same slide layout as the template.
    new_slide = ppt.slides.add_slide(slide.slide_layout)

    # Remove the automatically created placeholders from the new slide.
    # We loop over a *copy* of the list (using list(...)) because we are deleting items.
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # Copy each shape from the template slide into the new slide.
    # This uses deepcopy so that the full XML of each shape is duplicated.
    for shape in slide.shapes:
        # new_shape = copy.deepcopy(shape.element)
        new_shape = parse_xml(shape.element.xml)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    return new_slide

def delete_slide(ppt, index):
    """
    Safely delete a slide at zero-based `index` from `ppt` without triggering repair mode.
    """
    # Access the slide ID list from presentation XML
    pres_part = ppt.part
    sldIdLst = pres_part._element.sldIdLst
    # Target the specific slide ID element
    sldId = list(sldIdLst)[index]
    rId = sldId.rId
    # Remove the slide ID element from the list
    sldIdLst.remove(sldId)
    # Remove the slide part relationship
    pres_part.drop_rel(rId)
    return ppt


def replace_and_format_title(shapes_with_tables, replacement_text, replaced_text):
    """
    Replaces 'Title' in a shape's text with replacement_text and applies formatting:
    - Font: Univers for KPMG
    - Size: 14 pt
    - Color: #FFFFFF
    """
    for shape in shapes_with_tables:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text == replaced_text:
                        if pd.isna(replacement_text) or str(replacement_text).strip() == "":
                            cell.text = ""  # Keep cell blank
                        else:
                            cell.text = str(replacement_text)

                        # Apply formatting to all paragraphs and runs in the cell
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = "Univers for KPMG"
                                run.font.size = Pt(14)
                                run.font.color.rgb = RGBColor(255, 255, 255) #white


def update_table_texts(shapes_with_tables, replacing_text, review_text):
    """
    Replaces table text in table cells with provided values,
    and applies formatting:
    - Font: Univers for KPMG
    - Size: 11 pt
    - Color: #002060
    """
    for shape in shapes_with_tables:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text == replacing_text:
                        if pd.isna(review_text) or str(review_text).strip() == "":
                            cell.text = ""  # Keep cell blank
                        else:
                            cell.text = str(review_text)

                        # Apply formatting to all paragraphs and runs in the cell
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = "Univers for KPMG"
                                run.font.size = Pt(11)
                                run.font.color.rgb = RGBColor(0x00, 0x20, 0x60) #002060

def check_overflow_by_chars(cell, max_chars):
    """
    Checks if the cell text exceeds max_chars.
    Returns (overflow_flag, first_chars, remaining_chars)
    """
    text = cell.text
    if len(text) > max_chars:
        return True, text[:max_chars], text[max_chars:]
    return False, text, ""

def format_cell(cell, font_name="Univers for KPMG", font_size=11, font_color=RGBColor(0x00, 0x20, 0x60)):
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color

def main(tracker_path, ppt_path, output_folder_path):
    issue_tracker_raw_df = pd.read_excel(r"%s" % (tracker_path), header = 3, engine="openpyxl")
    ppt = Presentation(r"%s" % (ppt_path))

    # Detailed Observation Section
    default_template = ppt.slides[25]
    new_slide_idx = 25

    for _, review_row in issue_tracker_raw_df.iterrows():
        # Duplicating template slide
        new_slide = duplicate_slide(ppt, default_template)
         
        # Arrange Slide
        # Access the internal slide ID list
        xml_slides = ppt.slides._sldIdLst
        slides = list(xml_slides)
        # Remove the newly added slide from the end...
        xml_slides.remove(slides[-1])
        # ...and insert it right after the template slide (i.e. at position template_index+1)
        new_slide_idx = new_slide_idx + 1
        xml_slides.insert(new_slide_idx, slides[-1])

        # Creating 2 lists for all textbox shapes in the template slides
        shapes_with_textbox = [shape for shape in new_slide.shapes if shape.has_text_frame]
        shapes_with_tables = [shape for shape in new_slide.shapes if shape.has_table]

        # Changing title based on the excel row        
        replace_and_format_title(shapes_with_tables, str(review_row["Title"]).strip(), "Name of the Review")

        # Changing the text in background slide
        update_table_texts(shapes_with_tables,"Risk-text", str(review_row["Risk"]).strip())
        update_table_texts(shapes_with_tables,"Current State Control-text", str(review_row["Test Attributes / Current State Control"]).strip())
        update_table_texts(shapes_with_tables,"Issue Description-text", str(review_row["Issue Description"]).strip())
        update_table_texts(shapes_with_tables,"Recommendation-text", str(review_row["Recommendation"]).strip())
        update_table_texts(shapes_with_tables,"Management Response-text", str(review_row["Management Response (F.Y. 2024)"]).strip())
        update_table_texts(shapes_with_tables,"Responsibility & Timelines-text", str(review_row["Responsibility & Timelines"]).strip())

        # '''
        # Define cell locations and their overflow character limits
        cell_configs = [
            {"name": "risk", "cell": shapes_with_tables[1].table.rows[1].cells[0], "max_chars": 450},
            {"name": "csc", "cell": shapes_with_tables[1].table.rows[3].cells[0], "max_chars": 500},
            {"name": "id", "cell": shapes_with_tables[1].table.rows[6].cells[0], "max_chars": 600},
            {"name": "rec", "cell": shapes_with_tables[1].table.rows[8].cells[0], "max_chars": 700},
            {"name": "mr", "cell": shapes_with_tables[1].table.rows[2].cells[1], "max_chars": 900},
            {"name": "resp", "cell": shapes_with_tables[1].table.rows[8].cells[1], "max_chars": 250}
        ]

        overflow_flag = False
        overflow_info = {}

        # Check each cell for overflow and store results
        for config in cell_configs:
            flag, first, remaining = check_overflow_by_chars(config["cell"], config["max_chars"])
            if flag:
                overflow_flag = True
                overflow_info[config["name"]] = {"first": first, "remaining": "-"+remaining}


        if overflow_flag:
            # Duplicate the slide and adjust its position
            new_slide1 = duplicate_slide(ppt, new_slide)
            xml_slides = ppt.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            new_slide_idx += 1
            xml_slides.insert(new_slide_idx, slides[-1])

            # Get shapes for both slides
            shapes_with_textbox1 = [shape for shape in new_slide1.shapes if shape.has_text_frame]
            shapes_with_tables1 = [shape for shape in new_slide1.shapes if shape.has_table]

            # Update titles for both slides
            for idx, suffix in [(0, " (1/2)"), (0, " (2/2)")]:
                title_shape = shapes_with_tables[1].table.rows[0].cells[idx] if suffix == " (1/2)" else shapes_with_tables1[1].table.rows[0].cells[idx]
                title_text = title_shape.text
                text_frame = title_shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = title_text + suffix
                run.font.name = "Univers for KPMG"
                run.font.size = Pt(14)
                run.font.color.rgb = RGBColor(255,255,255)

            # Map cell names to their locations in both slides
            cell_map = {
                "risk":  {"orig": shapes_with_tables[1].table.rows[2].cells[0], "dup": shapes_with_tables1[1].table.rows[2].cells[0]},
                "csc": {"orig": shapes_with_tables[1].table.rows[4].cells[0], "dup": shapes_with_tables1[1].table.rows[4].cells[0]},
                "id": {"orig": shapes_with_tables[1].table.rows[6].cells[0], "dup": shapes_with_tables1[1].table.rows[6].cells[0]},
                "rec": {"orig": shapes_with_tables[1].table.rows[8].cells[0], "dup": shapes_with_tables1[1].table.rows[8].cells[0]},
                "mr": {"orig": shapes_with_tables[1].table.rows[2].cells[1], "dup": shapes_with_tables1[1].table.rows[2].cells[1]},
                "resp": {"orig": shapes_with_tables[1].table.rows[8].cells[1], "dup": shapes_with_tables1[1].table.rows[8].cells[1]}
                }

            # Update cell texts and formatting for overflowed cells
            for key, info in overflow_info.items():
                # Original slide: show first part
                cell_map[key]["orig"].text = info["first"]
                format_cell(cell_map[key]["orig"])
                # Duplicated slide: show remaining part
                cell_map[key]["dup"].text = info["remaining"]
                format_cell(cell_map[key]["dup"]) # '''
    
    # Deleting the original template slide
    ppt = delete_slide(ppt, 25)
    # Save the final presentation
    
    # Define the output folder inside the parent directory
    output_folder_path = Path(output_folder_path)
    output_folder = output_folder_path / "ICOFR Report"
    # Create the folder if it doesn't exist
    output_folder.mkdir(parents=True, exist_ok=True)

    OP_DIR = output_folder
    if not os.path.exists(OP_DIR):
            os.makedirs(OP_DIR)
    output_path = r"%s/DIB - ICOFR 2024 - Detailed Gap Report - Template_Automated.pptx" % (OP_DIR)
    ppt.save(output_path)

main(
    r"C:\Users\aryansharma8\OneDrive - KPMG\Documents\POC\FS Automation\ICOFR\DIB - ICOFR 2024 - Detailed Gap Report- Automation Requirement.xlsx",
    r"C:\Users\aryansharma8\OneDrive - KPMG\Documents\POC\FS Automation\ICOFR\DIB - ICOFR 2024 - Detailed Gap Report - Template.pptx",
    r"C:\Users\aryansharma8\OneDrive - KPMG\Desktop\FS Utility Output"
    )
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from datetime import datetime
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
import copy
from copy import deepcopy
from pptx.oxml.ns import qn
import tkinter as tk
from tkinter import filedialog, messagebox
import os
from pathlib import Path
import shutil
import subprocess
from PIL import Image, ImageTk
import time

# Function to duplicate a slide by copying its XML element tree
def duplicate_slide(ppt, slide):
    # Create a new slide using the same slide layout as the template.
    new_slide = ppt.slides.add_slide(slide.slide_layout)

    # Remove the automatically created placeholders from the new slide.
    # We loop over a *copy* of the list (using list(...)) because we are deleting items.
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # Copy each shape from the template slide into the new slide.
    # This uses deepcopy so that the full XML of each shape is duplicated.
    for shape in slide.shapes:
        # new_shape = copy.deepcopy(shape.element)
        new_shape = parse_xml(shape.element.xml)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    return new_slide


def delete_slide(ppt, index):
    """
    Safely delete a slide at zero-based `index` from `ppt` without triggering repair mode.
    """
    # Access the slide ID list from presentation XML
    pres_part = ppt.part
    sldIdLst = pres_part._element.sldIdLst
    # Target the specific slide ID element
    sldId = list(sldIdLst)[index]
    rId = sldId.rId
    # Remove the slide ID element from the list
    sldIdLst.remove(sldId)
    # Remove the slide part relationship
    pres_part.drop_rel(rId)
    return ppt


def delete_extra_rows(table, keep_row_count):
    """
    Remove all rows in the PPTX table beyond `keep_row_count`.
    `keep_row_count` should include header rows if you have them.
    """
    tbl_elm = table._tbl           # the <a:tbl> element
    all_rows = list(tbl_elm.tr_lst)  # list of <a:tr> elements
    for row_elm in all_rows[keep_row_count:]:
        tbl_elm.remove(row_elm)

def update_page_no(ppt, slide_idx, reference_text, replace_text):
    # Find the index of the slide containing the reference text
    part1_idx = None
    for idx, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text == reference_text:
                part1_idx = idx
                break
        if part1_idx is not None:
            break

    if part1_idx is None:
        return  # Reference text not found

    # Update the target slide's table cell
    target_slide = ppt.slides[slide_idx]
    for shape in target_slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text == replace_text:
                        cell.text = str(part1_idx + 1)
                        # Format the cell text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(14)

def Concurrent_Report_Main(tracker_path, ppt_path, report_type, output_folder_path, month, year):
    issue_tracker_raw_df = pd.read_excel(r"%s" % (tracker_path), header = 1)
    status_list = issue_tracker_raw_df['Draft/Updated Draf/Final'].unique().tolist()
    # Convert input month name to month number
    input_month_num = datetime.strptime(month, "%B").month
    input_year = int(year)
    df_dashboard = pd.read_excel(tracker_path, sheet_name="Volumes", engine="openpyxl")

    for status in status_list:
        if report_type == "Both":
            pass
        elif status != report_type:
            continue
        status_rept_df = issue_tracker_raw_df[issue_tracker_raw_df['Draft/Updated Draf/Final'] == status]
        reports_list = status_rept_df['Name of the report'].unique().tolist()
        past_issues_df = status_rept_df.copy(deep=True)
        # Convert 'Month-Year' column to datetime format
        past_issues_df['Month-Year'] = pd.to_datetime(past_issues_df['Month-Year'], errors='coerce')

        # Filter rows where 'Month-Year' is less than the input month and year
        final_past_issues_df = past_issues_df[
            (past_issues_df['Status'] == 'Open') &
            ((past_issues_df['Month-Year'].dt.year < input_year) |
             ((past_issues_df['Month-Year'].dt.year == input_year) &
              (past_issues_df['Month-Year'].dt.month < input_month_num)))
        ]

        # Filter rows where 'Month-Year' is greater than or equal to the input month and year
        audit_mnt_issues_df = past_issues_df[
            ((past_issues_df['Month-Year'].dt.year > input_year) |
             ((past_issues_df['Month-Year'].dt.year == input_year) &
              (past_issues_df['Month-Year'].dt.month == input_month_num)))
        ]


        # Report level iterration
        for report in reports_list:
            # print("\nCurrent Report is: ", report )
            # # print(report)
            pi_slide_idx = 15
            audit_mnt_slide_idx = 12
            exec1_slide_idx = 7
            exec2_slide_idx = 8

            ppt = Presentation(r"%s" % (ppt_path))
            # ppt = Presentation(r"Input_Data\ABC Bank Report_updated.pptx")
            first_slide = ppt.slides[0]
            default_template = ppt.slides[12]
            exec_summ_slide = ppt.slides[7]
            exec_summ_slide2 = ppt.slides[8]
            dashboard_slide = ppt.slides[4]
            dashboard_slide_idx = 4

            # Preprocessing issue tracker
            curr_rept_df = status_rept_df[status_rept_df['Name of the report'] == report].reset_index()
            total_obs_count = curr_rept_df.shape[0]
            # print("Total Observation count: ", total_obs_count)
            rep_final_past_issues_df = final_past_issues_df[
                final_past_issues_df['Name of the report'] == report].reset_index()
            pi_obs_count = rep_final_past_issues_df.shape[0]
            # print("Past Observation Count: ", pi_obs_count)
            rep_audit_mnt_issues_df = audit_mnt_issues_df[
                audit_mnt_issues_df['Name of the report'] == report].reset_index()
            aud_mnt_obs_count = rep_audit_mnt_issues_df.shape[0]
            # print("Detailed Observation Count: ", aud_mnt_obs_count)
            # curr_rept_df['Report Month-Year'] = pd.to_datetime(curr_rept_df['Report Month-Year'], errors='coerce')

            ## Changes in 1st slide
            shapes_with_textbox = [shape for shape in first_slide.shapes if shape.has_text_frame]
            if shapes_with_textbox[3].text.strip().lower() == 'b b s r & associates llp':
                curr_rept_df = curr_rept_df[
                    curr_rept_df['Entity'].str.strip().str.lower() == 'b b s r & associates llp'].reset_index()
            else:
                curr_rept_df = curr_rept_df[
                    curr_rept_df['Entity'].str.strip().str.lower() != 'b b s r & associates llp'].reset_index()

            # Update first slide
            shapes_with_textbox[0].text = status + ' Concurrent Audit Report for the month of %s' % (
                curr_rept_df.loc[0, 'Report Month-Year'].strftime("%B-%Y"))
            paragraph = shapes_with_textbox[0].text_frame.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = "Arial (Body)"
            run.font.size = Pt(14)


            # First page Report Name
            shapes_with_textbox[1].text = 'ABC Bank\n\n%s' % (report)
            text_frame1 = shapes_with_textbox[1].text_frame
            for idx, para in enumerate(text_frame1.paragraphs):
                # # print(para.text)
                for run in para.runs:
                    run.font.name = "KPMG Extralight (Headings)"  # “Aptos (Body)” → just "Aptos"
                    run.font.size = Pt(60)
                if ' ' in para.text:  # Check if the paragraph contains a space
                    for run in para.runs:
                        run.font.name = "KPMG Extralight (Headings)"
                        run.font.size = Pt(60)

            # rep_final_past_issues_df['Target Closure Date'] = pd.to_datetime(rep_final_past_issues_df['Target Closure Date'], errors='coerce')
            shapes_with_textbox[2].text = curr_rept_df.loc[0, 'Month of Issuance']
            paragraph = shapes_with_textbox[2].text_frame.paragraphs[0]
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = "Arial"
            run.font.size = Pt(12)

            # # Detail Observation Updation
            # # Observation slide level itteration
            cols_to_check = ['L', 'D', 'R', 'H', 'S', 'I', 'O/P', 'O ', 'M']

            def observations(curr_rept_df, obs_count, new_slide_idx, pi_flag):
                idx_month = 0
                for index in range(obs_count):
                    if pi_flag:
                        # print("Current pi index is: ",index + 1, "out of: ", obs_count)
                        pass
                    else:
                        # print("Current index is: ",index + 1, "out of: ", obs_count)  
                        pass                      
                    new_slide = duplicate_slide(ppt, default_template)

                    # Arrange Slide
                    # Access the internal slide ID list
                    xml_slides = ppt.slides._sldIdLst
                    slides = list(xml_slides)
                    # Remove the newly added slide from the end...
                    xml_slides.remove(slides[-1])
                    # ...and insert it right after the template slide (i.e. at position template_index+1)
                    new_slide_idx = new_slide_idx + 1
                    xml_slides.insert(new_slide_idx, slides[-1])

                    shapes_with_textbox = [shape for shape in new_slide.shapes if shape.has_text_frame]
                    shapes_with_tables = [shape for shape in new_slide.shapes if shape.has_table]

                    if pi_flag == True:
                        # Convert 'Report Month-Year' to datetime
                        curr_rept_df['Pi_Month'] = pd.to_datetime(curr_rept_df['Report Month-Year'], format='%b-%Y')
                        curr_rept_df['Pi_Year'] = pd.to_datetime(curr_rept_df['Report Month-Year'], format='%b-%Y')

                        # Extract year  #Aryan
                        curr_rept_df['Pi_Year'] = curr_rept_df['Pi_Month'].dt.year
                        # Extract month name
                        curr_rept_df['Pi_Month'] = curr_rept_df['Pi_Month'].dt.month_name()

                        time_period_month = curr_rept_df.loc[index, 'Pi_Month']
                        time_period_year = curr_rept_df.loc[index, 'Pi_Year']
                        current_month = curr_rept_df.loc[index, 'Pi_Month']

                        if (index < 1):  # Aryan
                            old_month = curr_rept_df.loc[index, 'Pi_Month']
                        else:
                            old_month = curr_rept_df.loc[index - 1, 'Pi_Month']

                        if (current_month == old_month):
                            idx_month = idx_month + 1
                        else:
                            idx_month = 1
                        # # print(index)
                        shapes_with_textbox[0].text = 'Detailed Observation - %s %s' % (
                        time_period_month, time_period_year)

                    else:
                        # Update title
                        # shapes_with_textbox[0].text = 'Detailed Observation (%s of %s)'%(index+1, total_obs_count)
                        shapes_with_textbox[0].text = 'Detailed Observation'

                    if pd.isna(curr_rept_df.loc[index, 'Root Cause ']):
                        # # print(shapes_with_tables[1].table.rows[3].cells[1].text)
                        shapes_with_tables[1].table.rows[1].cells[1].text = "NA"
                    else:
                        shapes_with_tables[1].table.rows[1].cells[1].text = str(curr_rept_df.loc[index, 'Root Cause '])
                    paragraph = shapes_with_tables[1].table.rows[1].cells[1].text_frame.paragraphs[0]
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    run.font.name = "Aptos (Body)"
                    run.font.size = Pt(14)

                    # Add text in MAP
                    if pd.isna(curr_rept_df.loc[index, 'MAP']):
                        # # print(shapes_with_tables[1].table.rows[3].cells[1].text)
                        shapes_with_tables[1].table.rows[3].cells[1].text = "NA"
                    else:
                        shapes_with_tables[1].table.rows[3].cells[1].text = str(curr_rept_df.loc[index, 'MAP'])
                    cell1 = shapes_with_tables[1].table.rows[3].cells[1]
                    for idx, para in enumerate(cell1.text_frame.paragraphs):
                        for run in para.runs:
                            run.font.name = "Aptos"  # “Aptos (Body)” → just "Aptos"
                            run.font.size = Pt(14)

                    # Add text in Responsibility
                    shapes_with_tables[1].table.rows[5].cells[1].text = curr_rept_df.loc[index, 'Responsibility']
                    # paragraph = shapes_with_tables[1].table.rows[5].cells[1].text_frame.paragraphs[0]
                    # run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    # run.font.name = "Aptos (Body)"
                    # run.font.size = Pt(14)
                    cell2 = shapes_with_tables[1].table.rows[5].cells[1]
                    for idx, para in enumerate(cell2.text_frame.paragraphs):
                        for run in para.runs:
                            run.font.name = "Aptos"  # “Aptos (Body)” → just "Aptos"
                            run.font.size = Pt(14)

                    # Add text in Target Closure Date
                    if pd.isna(curr_rept_df.loc[index, 'Target Closure Date']):
                        shapes_with_tables[1].table.rows[5].cells[2].text = "NA"
                    else:
                        shapes_with_tables[1].table.rows[5].cells[2].text = curr_rept_df.loc[
                            index, 'Target Closure Date']
                    # paragraph = shapes_with_tables[1].table.rows[5].cells[2].text_frame.paragraphs[0]
                    # run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    # run.font.name = "Aptos (Body)"
                    # run.font.size = Pt(14)
                    cell3 = shapes_with_tables[1].table.rows[5].cells[2]
                    for idx, para in enumerate(cell3.text_frame.paragraphs):
                        for run in para.runs:
                            run.font.name = "Aptos"  # “Aptos (Body)” → just "Aptos"
                            run.font.size = Pt(14)

                    # Create an empty list to store valid column indices
                    non_na_indices = []

                    # Check each column for non-NA at row index 0
                    for idx, col in enumerate(cols_to_check):  # zero-based index
                        if str(curr_rept_df.loc[index, col]).strip().lower() == 'yes':
                            non_na_indices.append(idx)

                    # non_na_indices = [8]
                    empty_text_index = 0

                    for shape in new_slide.shapes:
                        if shape.shape_type == 6:  # GROUP
                            for sub1 in shape.shapes:
                                if sub1.shape_type == 6:  # Nested GROUP
                                    for sub2 in sub1.shapes:
                                        if hasattr(sub2, "text") and sub2.has_text_frame:
                                            if sub2.text.strip() == "":
                                                if empty_text_index in non_na_indices:
                                                    sub2.text = u"\u2713"  # tick
                                                    # # print(f"Updated shape #{empty_text_index + 1} with tick.")
                                                    paragraph = sub2.text_frame.paragraphs[0]
                                                    paragraph.runs[0].font.size = Pt(14)
                                                    paragraph.alignment = PP_ALIGN.CENTER
                                                    sub2.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                                                empty_text_index += 1

                    # Add text in Observation Header & Observation Details
                    shapes_with_tables[1].table.rows[1].cells[0].text = (
                        shapes_with_tables[1].table.rows[1].cells[0].text).replace(
                        "Title of Observation/Issue description ( Column C)",
                        curr_rept_df.loc[index, 'Issue disc']).replace("Detailed Observation  (Column D)",
                                                                       curr_rept_df.loc[
                                                                           index, 'Detailed Observation'])

                    cell = shapes_with_tables[1].table.rows[1].cells[0]
                    # Check if the cell text contains more than 10 lines (using newline as the delimiter)
                    lines = cell.text.split("\n")
                    # Store first 10 and remaining lines in separate variables
                    first_ten_lines = "\n".join(lines[:10])
                    remaining_lines = "\n".join(lines[10:])

                    # Extract the first line from first_ten_lines
                    first_line = first_ten_lines.splitlines()[0]
                    if pi_flag:  # Aryan
                        first_line = "%s." % (idx_month) + first_line
                    else:
                        first_line = "%s." % (index + 1) + first_line
                    first_line = first_line.rstrip(".,;") + " (2 of 2)"

                    # Prepend the first line to the remaining_lines
                    remaining_lines = first_line + "\n\n" + remaining_lines

                    if len(lines) > 10:
                        # Modify the first line of 'first_ten_lines' with (1/2)
                        first_ten_lines_lines = first_ten_lines.splitlines()
                        if pi_flag:
                            first_ten_lines_lines[0] = "%s." % (idx_month) + first_ten_lines_lines[0]
                        else:
                            first_ten_lines_lines[0] = "%s." % (index + 1) + first_ten_lines_lines[0]
                        first_ten_lines_lines[0] = first_ten_lines_lines[0].rstrip(".,;") + " (1 of 2)"
                        shapes_with_tables[1].table.rows[1].cells[0].text = "\n".join(first_ten_lines_lines)
                    else:
                        first_ten_lines_lines2 = first_ten_lines.splitlines()
                        if pi_flag:
                            first_ten_lines_lines2[0] = "%s." % (idx_month) + first_ten_lines_lines2[0].rstrip(".,;")
                        else:
                            first_ten_lines_lines2[0] = "%s." % (index + 1) + first_ten_lines_lines2[0].rstrip(".,;")
                        shapes_with_tables[1].table.rows[1].cells[0].text = "\n".join(first_ten_lines_lines2)
                    cell = shapes_with_tables[1].table.rows[1].cells[0]

                    for idx, para in enumerate(cell.text_frame.paragraphs):
                        # # print(para.text)
                        for run in para.runs:
                            if idx == 0:
                                run.font.bold = True
                            run.font.name = "Aptos"  # “Aptos (Body)” → just "Aptos"
                            run.font.size = Pt(14)

                    if len(lines) > 10:
                        new_slide1 = duplicate_slide(ppt, new_slide)
                        # Access the internal slide ID list
                        xml_slides = ppt.slides._sldIdLst
                        slides = list(xml_slides)
                        # Remove the newly added slide from the end...
                        xml_slides.remove(slides[-1])
                        # ...and insert it right after the template slide (i.e. at position template_index+1)
                        new_slide_idx = new_slide_idx + 1
                        xml_slides.insert(new_slide_idx, slides[-1])

                        shapes_with_textbox = [shape for shape in new_slide1.shapes if shape.has_text_frame]
                        shapes_with_tables = [shape for shape in new_slide1.shapes if shape.has_table]

                        shapes_with_tables[1].table.rows[1].cells[0].text = remaining_lines
                        cell = shapes_with_tables[1].table.rows[1].cells[0]

                        for idx, para in enumerate(cell.text_frame.paragraphs):
                            # # print(para.text)
                            for run in para.runs:
                                if idx == 0:
                                    run.font.bold = True
                                run.font.name = "Aptos"  # “Aptos (Body)” → just "Aptos"
                                run.font.size = Pt(14)

            observations(rep_final_past_issues_df, pi_obs_count, pi_slide_idx, True)
            observations(rep_audit_mnt_issues_df, aud_mnt_obs_count, audit_mnt_slide_idx, False)
            # delete_slide(ppt, 12)
            # delete_slide(ppt, audit_mnt_slide_idx)

            ## Exec Summ2 (Details of Observations)
            for index in range(total_obs_count):
                if index % 6 == 0:
                    exe2_new_slide = duplicate_slide(ppt, exec_summ_slide2)
                    exec2_summ_tables = [shape for shape in exe2_new_slide.shapes if shape.has_table]
                    count2_exec_slide2_row = 0

                    # Arrange Slide
                    # Access the internal slide ID list
                    xml_slides = ppt.slides._sldIdLst
                    slides = list(xml_slides)
                    # Remove the newly added slide from the end...
                    xml_slides.remove(slides[-1])
                    # ...and insert it right after the template slide (i.e. at position template_index+1)
                    exec2_slide_idx = exec2_slide_idx + 1
                    xml_slides.insert(exec2_slide_idx, slides[-1])

                count2_exec_slide2_row = count2_exec_slide2_row + 1
                exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[0].text = str(index + 1)
                paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[0].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

                exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[1].text = curr_rept_df.loc[
                    index, 'Issue disc']
                paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[1].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)

                exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[4].text = curr_rept_df.loc[
                    index, 'Repeat Observation']
                paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[4].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

                exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[5].text = str(
                    curr_rept_df.loc[index, 'No of Instances'])
                paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[5].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

                # Create an empty list to store valid column indices
                non_na_indices = []
                # Check each column for non-NA at row index 0
                for idx, col in enumerate(cols_to_check):  # zero-based index
                    if str(curr_rept_df.loc[index, col]).strip().lower() == 'yes':
                        non_na_indices.append(idx)
                print(non_na_indices)

                # cols_to_check = ['L', 'D', 'R', 'H', 'S', 'I', 'O/P', 'O ', 'M'] [2, 5, 6, 7, 8]
                for val in non_na_indices:
                    cell = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[2]
                    cell.fill.solid()
                    if val == 0:
                        exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[2].text = "Low"
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 0)
                    if val == 3:
                        exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[2].text = "High"
                        cell.fill.fore_color.rgb = RGBColor(255, 0, 0)
                    if val == 8:
                        exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[2].text = "Medium"
                        cell.fill.fore_color.rgb = RGBColor(255, 165, 0)
                    cell = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[2]
                    paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[2].text_frame.paragraphs[0]
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    run.font.name = "Aptos (Body)"
                    run.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.CENTER 

                    cell3_loc = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[3]
                    if val == 1:
                        if cell3_loc.text.strip() != "":
                            cell3_loc.text = cell3_loc.text + ','
                        cell3_loc = cell3_loc.text + 'D'
                    if val == 7:
                        if cell3_loc.text.strip() != "":
                            cell3_loc.text = cell3_loc.text + ','
                        cell3_loc.text = cell3_loc.text + 'O'
                    if val == 4:
                        if cell3_loc.text.strip() != "":
                            cell3_loc.text = cell3_loc.text + ','
                        cell3_loc.text = cell3_loc.text + 'S'

                    paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[3].text_frame.paragraphs[
                        0]
                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                    run.font.name = "Aptos (Body)"
                    run.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.CENTER

                    cell6_loc = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[6]
                    if val == 2:
                        if cell6_loc.text.strip() != "":
                            cell6_loc.text = cell6_loc.text + ','
                        cell6_loc.text = cell6_loc.text + 'R'
                    if val == 6:
                        if cell6_loc.text.strip() != "":
                            cell6_loc.text = cell6_loc.text + ','
                        cell6_loc.text = cell6_loc.text + 'O/P'
                    if val == 5:
                        if cell6_loc.text.strip() != "":
                            cell6_loc.text = cell6_loc.text + ','
                        cell6_loc.text = cell6_loc.text + 'I'

                paragraph = exec2_summ_tables[0].table.rows[count2_exec_slide2_row].cells[6].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

            # Delete extra rows from table
            table_shape = exec2_summ_tables[0]
            tbl = table_shape.table
            filled_rows = count2_exec_slide2_row + 1
            delete_extra_rows(tbl, filled_rows)
            # delete_slide(ppt, 8)

            ## Exe Summ-1 (Summary of Observations)
            # print("Total Observation Count: ",total_obs_count)
            for index1 in range(total_obs_count):

                if index1 % 2 == 0:
                    if index1 == 0:
                        exe_new_slide = duplicate_slide(ppt, exec_summ_slide)
                    else:
                        exe_new_slide = duplicate_slide(ppt, exe_new_slide)
                    exec_summ_tables = [shape for shape in exe_new_slide.shapes if shape.has_table]
                    count_exec_slide_row = 0
                    # # print("index for slide created %s"%(index1))
                    # # print("table for slide created %s"%(exec_summ_tables))

                    # Arrange Slide
                    # Access the internal slide ID list
                    xml_slides = ppt.slides._sldIdLst
                    slides = list(xml_slides)
                    # Remove the newly added slide from the end...
                    # # print("Slide ids %s"%slides)
                    xml_slides.remove(slides[-1])
                    # ...and insert it right after the template slide (i.e. at position template_index1+1)
                    exec1_slide_idx = exec1_slide_idx + 1
                    xml_slides.insert(exec1_slide_idx, slides[-1])

                count_exec_slide_row = count_exec_slide_row + 1
                exec_summ_tables[0].table.rows[count_exec_slide_row].cells[0].text = str(index1 + 1)
                paragraph = exec_summ_tables[0].table.rows[count_exec_slide_row].cells[0].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

                exec_summ_tables[0].table.rows[count_exec_slide_row].cells[1].text = curr_rept_df.loc[
                    index1, 'Issue disc']
                paragraph = exec_summ_tables[0].table.rows[count_exec_slide_row].cells[1].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)

                exec_summ_tables[0].table.rows[count_exec_slide_row].cells[2].text = str(
                    curr_rept_df.loc[index1, 'No of Instances'])
                paragraph = exec_summ_tables[0].table.rows[count_exec_slide_row].cells[2].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

                exec_summ_tables[0].table.rows[count_exec_slide_row].cells[3].text = str(
                    curr_rept_df.loc[index1, 'Repeat Observation'])
                paragraph = exec_summ_tables[0].table.rows[count_exec_slide_row].cells[3].text_frame.paragraphs[0]
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                run.font.name = "Aptos (Body)"
                run.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER


                 # Add text in MAP
                if pd.isna(curr_rept_df.loc[index1, 'MAP']):
                    # # print(exec_summ_tables[1].table.rows[count_exec_slide_row].cells[1].text)
                    exec_summ_tables[0].table.rows[count_exec_slide_row].cells[4].text = "NA"
                else:
                    exec_summ_tables[0].table.rows[count_exec_slide_row].cells[4].text = str(curr_rept_df.loc[index1, 'MAP'])
                cell1 = exec_summ_tables[0].table.rows[count_exec_slide_row].cells[4]
                for idx, para in enumerate(cell1.text_frame.paragraphs):
                    for run in para.runs:
                        run.font.name = "Aptos"  # “Aptos (Body)” → just "Aptos"
                        run.font.size = Pt(14)
            #Delete extra rows from table
            table_shape = exec_summ_tables[0]
            tbl = table_shape.table
            filled_rows = count_exec_slide_row + 1
            delete_extra_rows(tbl, filled_rows)

            # delete_slide(ppt, 7)
            cal_exec2_slide_idx = exec1_slide_idx + 1
            cal_exec1_slide_idx = 7
            cal_obs_slide_idx = cal_exec2_slide_idx + (exec2_slide_idx - 8) + 4
            # print(f"Initial slide count: {len(ppt.slides)}")
            # Indices of slides to delete (0-based)
            slides_to_delete = [cal_exec2_slide_idx, cal_exec1_slide_idx, cal_obs_slide_idx]

            # # Delete in reverse order to avoid shifting indices
            # slides_to_delete = [7, 14, 20] # Example: Delete slides at index 0, 2, and 4
            to_delete = sorted(slides_to_delete, reverse=True)

            # Delete the slides
            # modified_ppt = create_presentation_excluding_slides(ppt, slides_to_delete)
            for idx in to_delete:
                ppt = delete_slide(ppt, idx)
            # print(f"Slide count after deletion: {len(ppt.slides)}")

            # '''# Changes in dashboard slide
            filtered_df = df_dashboard[df_dashboard["Name of the report"] == report].copy()
            filtered_df["Month-Year"] = pd.to_datetime(filtered_df["Month-Year"]).dt.strftime("%b")

            # Pivot and transpose data
            pivot_df = filtered_df.pivot_table(index="Month-Year", columns="Area", values="Volume", aggfunc="sum").fillna(0)
            transposed_df = pivot_df.T  # Rows = Area, Columns = Month-Year
            areas = transposed_df.index.tolist()

            """
            if len(areas) > 4 and len(areas) <= 8:
                # Split areas into two halves
                mid = len(areas) // 2
                area_groups = [areas[:mid], areas[mid:]]

                dashboard_slide_dup = duplicate_slide(ppt, dashboard_slide)
                # Access the internal slide ID list
                xml_slides = ppt.slides._sldIdLst
                slides = list(xml_slides)
                # Remove the newly added slide from the end...
                xml_slides.remove(slides[-1])
                # ...and insert it right after the template slide (i.e. at position template_index+1)
                dashboard_slide_idx = dashboard_slide_idx + 1
                xml_slides.insert(dashboard_slide_idx, slides[-1])
                slides = [dashboard_slide, dashboard_slide_dup]

            elif len(areas) > 8: 

                # Split into 3 groups
                split_size = len(areas) // 3
                area_groups = [
                    areas[:split_size],
                    areas[split_size:2*split_size],
                    areas[2*split_size:]
                ]

                dashboard_slide_dup = duplicate_slide(ppt, dashboard_slide)
                # Reorder slides
                xml_slides = ppt.slides._sldIdLst
                slides_xml = list(xml_slides)

                # Remove newly added slides from end
                xml_slides.remove(slides_xml[-1])

                # Insert them after the original slide
                dashboard_slide_idx += 1
                xml_slides.insert(dashboard_slide_idx, slides_xml[-1])  # dup2

                dashboard_slide_dup1 = duplicate_slide(ppt, dashboard_slide)
                
            # Reorder slides
                xml_slides = ppt.slides._sldIdLst
                slides_xml = list(xml_slides)

                # Remove newly added slides from end
                xml_slides.remove(slides_xml[-1])

                # Insert them after the original slide
                dashboard_slide_idx += 1
                xml_slides.insert(dashboard_slide_idx, slides_xml[-1])  # dup2

                # Final slide list
                slides = [dashboard_slide, ppt.slides[dashboard_slide_idx - 1], ppt.slides[dashboard_slide_idx]]

            else:
                area_groups = [areas]
                slides = [dashboard_slide]
                """
            
            area_groups = [areas]
            slides = [dashboard_slide]


            for slide, area_subset in zip(slides, area_groups):
                # Filter transposed_df for the current subset of areas
                subset_df = transposed_df.loc[area_subset]

                # Locate and remove existing chart
                chart_shape = next(shape for shape in slide.shapes if shape.has_chart)
                x, y, cx, cy = chart_shape.left, chart_shape.top, chart_shape.width, chart_shape.height
                slide.shapes._spTree.remove(chart_shape._element)

                # Create new chart data
                chart_data = CategoryChartData()
                chart_data.categories = list(subset_df.index)

                for month in subset_df.columns:
                    chart_data.add_series(month, list(subset_df.loc[:, month]))

                # Add new chart
                chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart

                # Customize chart
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False

                chart.has_title = True
                chart.chart_title.text_frame.text = "Volumes by Month and Area"
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

                chart.category_axis.tick_labels.font.size = Pt(10)
                chart.value_axis.tick_labels.font.size = Pt(10)
                chart.value_axis.has_major_gridlines = False

                for series in chart.series:
                    series.has_data_labels = True
                    data_labels = series.data_labels
                    data_labels.show_value = True
                    data_labels.font.size = Pt(10)
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END #'''
            
            # Page Number section
            update_page_no(ppt, 1,"Part I: Scope and Approach","a")
            update_page_no(ppt, 1,"Part II: Executive Summary","b")
            update_page_no(ppt, 1,"Part III: Detailed Observations","c")
            update_page_no(ppt, 1,"Part IV: Key statement of Facts","d")
            update_page_no(ppt, 1,"Part V: Past Issues Pending Resolution","e")            
            update_page_no(ppt, 1,"Part VI: Annexures","f")
            update_page_no(ppt, 1,"Part VII: Disclaimer and Note to Reader","g")            

            # Save the presentation
            output_folder_path = Path(output_folder_path)
            output_folder = output_folder_path / "Concurrent Report"
            # Create the folder if it doesn't exist
            output_folder.mkdir(parents=True, exist_ok=True)

            OP_DIR = output_folder / status
            if not os.path.exists(OP_DIR):
                    os.makedirs(OP_DIR)
            output_path = r"%s/%s_Report.pptx" % (OP_DIR, report)
            ppt.save(output_path)

    if report_type == "Both":
        complete_path = os.path.abspath(os.path.dirname(OP_DIR))
    else:
        complete_path = os.path.abspath(OP_DIR)

    return complete_path

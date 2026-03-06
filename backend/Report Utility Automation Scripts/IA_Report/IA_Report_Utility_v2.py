import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
import copy
from copy import deepcopy
from pptx.oxml.ns import qn
from openpyxl import load_workbook
import tkinter as tk
from tkinter import filedialog, messagebox
import os
from pathlib import Path
import shutil
import subprocess
from PIL import Image, ImageTk
import time

# ------Utility Functions-------
# ------IA Report Utility-------

# Function to duplicate a slide by copying its XML element tree
def duplicate_slide(ppt, slide):
    # Create a new slide using the same slide layout as the template.
    new_slide = ppt.slides.add_slide(slide.slide_layout)

    # Remove the automatically created placeholders from the new slide.
    # We loop over a *copy* of the list (using list(...)) because we are deleting items.
    for shape in list(new_slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)

    # Copy each shape from the template slide into the new slide.
    # This uses deepcopy so that the full XML of each shape is duplicated.
    for shape in slide.shapes:
        # new_shape = copy.deepcopy(shape.element)
        new_shape = parse_xml(shape.element.xml)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')

    return new_slide

def delete_slide(ppt, index):
    """
    Safely delete a slide at zero-based `index` from `ppt` without triggering repair mode.
    """
    # Access the slide ID list from presentation XML
    pres_part = ppt.part
    sldIdLst = pres_part._element.sldIdLst
    # Target the specific slide ID element
    sldId = list(sldIdLst)[index]
    rId = sldId.rId
    # Remove the slide ID element from the list
    sldIdLst.remove(sldId)
    # Remove the slide part relationship
    pres_part.drop_rel(rId)
    return ppt

def replace_and_format_title(shape, replacement_text):
    """
    Replaces 'Title' in a shape's text with replacement_text and applies formatting:
    - Font: Univers for KPMG
    - Size: 18 pt
    - Color: #002060
    """
    if shape.has_text_frame and "Title" in shape.text:
        text_frame = shape.text_frame
        text_frame.clear()

        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = replacement_text
        run.font.name = "Univers for KPMG"
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(0x00, 0x20, 0x60)  # #002060

def update_table_texts(shapes_with_tables, replacing_text, review_text):
    """
    Replaces table text in table cells with provided values,
    and applies formatting:
    - Font: Univers for KPMG
    - Size: 11 pt
    - Color: #002060
    """
    for shape in shapes_with_tables:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text == replacing_text:
                        if pd.isna(review_text) or str(review_text).strip() == "":
                            cell.text = ""  # Keep cell blank
                        else:
                            cell.text = str(review_text)

                        # Apply formatting to all paragraphs and runs in the cell
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.name = "Univers for KPMG"
                                run.font.size = Pt(11)
                                run.font.color.rgb = RGBColor(0x00, 0x20, 0x60) #002060

def update_risk_rating(shapes_with_textbox, RR):  
    for shape in shapes_with_textbox:
        if shape.text == 'RR':                
            fill = shape.fill
            fill.solid()  # Set fill type to solid
            if RR == 'H':
                fill.fore_color.rgb = RGBColor(0xFF, 0x00, 0x00)  # Red color
                shape.text = 'H'
            elif RR == 'M':
                fill.fore_color.rgb = RGBColor(255, 192, 0x00) # Yellow color
                shape.text = 'M'
            elif RR == 'L':
                fill.fore_color.rgb = RGBColor(146, 208, 80) # Green color
                shape.text = 'L'
            
            # Apply text formatting            
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text_frame.clear()  # Clear existing text
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER  # Center align the text
                run = p.add_run()
                run.text = RR
                run.font.name = "Univers for KPMG"
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)  # White

def update_risk_table(shapes_with_tables, replacing_text, review_text):
    for shape in shapes_with_tables:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text == replacing_text:
                        if review_text == "Yes":
                            fill = cell.fill
                            fill.solid()  # Set fill type to solid
                            fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x60)

def delete_extra_rows(table, keep_row_count):
    """
    Remove all rows in the PPTX table beyond `keep_row_count`.
    `keep_row_count` should include header rows if you have them.
    """
    tbl_elm = table._tbl           # the <a:tbl> element
    all_rows = list(tbl_elm.tr_lst)  # list of <a:tr> elements
    for row_elm in all_rows[keep_row_count:]:
        tbl_elm.remove(row_elm)

def update_page_no(ppt, slide_idx, reference_text, replace_text):
    # Find the index of the slide containing the reference text
    part1_idx = None
    for idx, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text == reference_text:
                part1_idx = idx
                break
        if part1_idx is not None:
            break

    if part1_idx is None:
        return  # Reference text not found

    # Update the target slide's table cell
    target_slide = ppt.slides[slide_idx]
    for shape in target_slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text == replace_text:
                        cell.text = str(part1_idx + 1)
                        # Format the cell text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.alignment = PP_ALIGN.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(14)

def check_overflow(cell, max_lines):
    """
    Checks if the cell text exceeds max_lines.
    Returns (overflow_flag, first_lines, remaining_lines)
    """
    lines = cell.text.split("\n")
    if len(lines) > max_lines:
        return True, "\n".join(lines[:max_lines]), "\n".join(lines[max_lines:])
    return False, cell.text, ""

def format_cell(cell, font_name="Univers for KPMG", font_size=11, font_color=RGBColor(0x00, 0x20, 0x60)):
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color

# def split_and_duplicate_slides(
#     ppt, 
#     original_slide, 
#     shapes_with_tables, 
#     shapes_with_textbox, 
#     new_slide_idx,
#     cell_limits = {"bg": 3, "obv": 8, "res": 8}
# ):
#     from pptx.util import Pt
#     from pptx.dml.color import RGBColor

#     def split_lines(text, max_lines):
#         lines = text.split("\n")
#         if len(lines) > max_lines:
#             return "\n".join(lines[:max_lines]), "\n".join(lines[max_lines:])
#         else:
#             return text, ""

#     def format_cell(cell):
#         for para in cell.text_frame.paragraphs:
#             for run in para.runs:
#                 run.font.name = "Univers for KPMG"
#                 run.font.size = Pt(11)
#                 run.font.color.rgb = RGBColor(0x00, 0x20, 0x60)

#     def update_title(slide, idx, total):
#         if total == 1:
#             return
#         title_shape = [shape for shape in slide.shapes if shape.has_text_frame][4]
#         title_text = title_shape.text
#         text_frame = title_shape.text_frame
#         text_frame.clear()
#         p = text_frame.paragraphs[0]
#         run = p.add_run()
#         run.text = f"{title_text} ({idx+1}/{total})"
#         run.font.name = "Univers for KPMG"
#         run.font.size = Pt(18)
#         run.font.color.rgb = RGBColor(0x00, 0x20, 0x60)

#     def get_cell(slide, key):
#         table = [shape for shape in slide.shapes if shape.has_table][0].table
#         if key == "bg":
#             return table.rows[1].cells[0]
#         elif key == "obv":
#             return table.rows[5].cells[1]
#         elif key == "res":
#             return table.rows[3].cells[0]
#         else:
#             raise ValueError("Unknown cell key")

#     cell_texts = {
#         "bg": shapes_with_tables[0].table.rows[1].cells[0].text,
#         "obv": shapes_with_tables[0].table.rows[5].cells[1].text,
#         "res": shapes_with_tables[0].table.rows[3].cells[0].text,
#     }

#     slides = [original_slide]
#     texts_remaining = cell_texts.copy()

#     while True:
#         slide = slides[-1]
#         for key in texts_remaining:
#             value = texts_remaining[key]
#             if isinstance(value, tuple):
#                 value = value[1]
#             first, remaining = split_lines(value, cell_limits[key])
#             cell = get_cell(slide, key)
#             cell.text = first
#             format_cell(cell)
#             texts_remaining[key] = remaining
#         if any(texts_remaining[key] for key in texts_remaining):
#             # Duplicate the current slide
#             next_slide = duplicate_slide(ppt, slide)
#             # Move the new slide right after the current slide
#             xml_slides = ppt.slides._sldIdLst
#             slides_list = list(xml_slides)
#             # Remove the new slide from the end
#             xml_slides.remove(slides_list[-1])
#             # Find the index of the current slide in the slide list
#             current_slide_id = slide.slide_id
#             for idx, sld in enumerate(slides_list):
#                 if sld.id == current_slide_id:
#                     insert_idx = idx + 1
#                     break
#             # Insert the new slide right after the current slide
#             xml_slides.insert(insert_idx, slides_list[-1])
#             slides.append(next_slide)
#         else:
#             break

#     total_slides = len(slides)
#     for idx, slide in enumerate(slides):
#         update_title(slide, idx, total_slides)

#     return slides

# ----Main Functions----   
def IA_Report_Utility_Main(tracker_path, ppt_path, output_folder_path):
    issue_tracker_raw_df = pd.read_excel(r"%s" % (tracker_path), header = 3, engine="openpyxl")
    ppt = Presentation(r"%s" % (ppt_path))

    # Detailed Observation Section
    default_template = ppt.slides[6]
    new_slide_idx = 6
    for _, review_row in issue_tracker_raw_df.iterrows():
        # Duplicating template slide
        new_slide = duplicate_slide(ppt, default_template)
         
        # Arrange Slide
        # Access the internal slide ID list
        xml_slides = ppt.slides._sldIdLst
        slides = list(xml_slides)
        # Remove the newly added slide from the end...
        xml_slides.remove(slides[-1])
        # ...and insert it right after the template slide (i.e. at position template_index+1)
        new_slide_idx = new_slide_idx + 1
        xml_slides.insert(new_slide_idx, slides[-1])

        # Creating 2 lists for all textbox shapes in the template slides
        shapes_with_textbox = [shape for shape in new_slide.shapes if shape.has_text_frame]
        shapes_with_tables = [shape for shape in new_slide.shapes if shape.has_table]

        # Changing title based on the excel row 
        for shape in shapes_with_textbox:
            replace_and_format_title(shape, str(review_row["Title"]))
        
        # Changing the text in background slide
        update_table_texts(shapes_with_tables,"Background-text", review_row["Background"])
        update_table_texts(shapes_with_tables,"Observation-text", review_row["Observation"])
        update_table_texts(shapes_with_tables,"Management Response-text", review_row["Management Response"])
        update_table_texts(shapes_with_tables,"Root Cause-text", review_row["Root Cause"])
        update_table_texts(shapes_with_tables,"Potential Implication-text", review_row["Implication"])
        update_table_texts(shapes_with_tables,"Recommendation-text", review_row["Recommendation"])
        update_table_texts(shapes_with_tables,"Due Date-text", review_row["Timeline"])
        update_table_texts(shapes_with_tables, "Responsibility-text", review_row["Responsibility"])

        # Change the Risk Rating
        update_risk_rating(shapes_with_textbox, review_row["Risk Rating"])

        # Update risk table
        update_risk_table(shapes_with_tables,"C",review_row["C"])
        update_risk_table(shapes_with_tables,"F",review_row["F"])
        update_risk_table(shapes_with_tables,"O",review_row["O"])
        update_risk_table(shapes_with_tables,"D",review_row["D"])
        update_risk_table(shapes_with_tables,"S",review_row["S"])
        
        # '''
        # Define cell locations and their overflow limits
        cell_configs = [
            {"name": "bg", "cell": shapes_with_tables[0].table.rows[1].cells[0], "max_lines": 3},
            {"name": "obv", "cell": shapes_with_tables[0].table.rows[5].cells[1], "max_lines": 8},
            {"name": "res", "cell": shapes_with_tables[0].table.rows[3].cells[0], "max_lines": 8},
            {"name": "rc", "cell": shapes_with_tables[0].table.rows[2].cells[1], "max_lines": 3},
            {"name": "pi", "cell": shapes_with_tables[0].table.rows[3].cells[2], "max_lines": 3},
            {"name": "rec", "cell": shapes_with_tables[0].table.rows[6].cells[0], "max_lines": 5},
            {"name": "resp", "cell": shapes_with_tables[0].table.rows[9].cells[1], "max_lines": 2},
            {"name": "due", "cell": shapes_with_tables[0].table.rows[9].cells[3], "max_lines": 2}
        ]

        overflow_flag = False
        overflow_info = {}

        # Check each cell for overflow and store results
        for config in cell_configs:
            flag, first, remaining = check_overflow(config["cell"], config["max_lines"])
            if flag:
                overflow_flag = True
                overflow_info[config["name"]] = {"first": first, "remaining": remaining}

        if overflow_flag:
            # Duplicate the slide and adjust its position
            new_slide1 = duplicate_slide(ppt, new_slide)
            xml_slides = ppt.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            new_slide_idx += 1
            xml_slides.insert(new_slide_idx, slides[-1])

            # Get shapes for both slides
            shapes_with_textbox1 = [shape for shape in new_slide1.shapes if shape.has_text_frame]
            shapes_with_tables1 = [shape for shape in new_slide1.shapes if shape.has_table]

            # Update titles for both slides
            for idx, suffix in [(4, " (1/2)"), (4, " (2/2)")]:
                title_shape = shapes_with_textbox[idx] if suffix == " (1/2)" else shapes_with_textbox1[idx]
                title_text = title_shape.text
                text_frame = title_shape.text_frame
                text_frame.clear()
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = title_text + suffix
                run.font.name = "Univers for KPMG"
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x00, 0x20, 0x60)

            # Map cell names to their locations in both slides
            cell_map = {
                "bg":  {"orig": shapes_with_tables[0].table.rows[1].cells[0], "dup": shapes_with_tables1[0].table.rows[1].cells[0]},
                "obv": {"orig": shapes_with_tables[0].table.rows[5].cells[1], "dup": shapes_with_tables1[0].table.rows[5].cells[1]},
                "res": {"orig": shapes_with_tables[0].table.rows[3].cells[0], "dup": shapes_with_tables1[0].table.rows[3].cells[0]},
                "rc": {"orig": shapes_with_tables[0].table.rows[2].cells[1], "dup": shapes_with_tables1[0].table.rows[2].cells[1]},
                "pi": {"orig": shapes_with_tables[0].table.rows[3].cells[2], "dup": shapes_with_tables1[0].table.rows[3].cells[2]},
                "rec": {"orig": shapes_with_tables[0].table.rows[6].cells[0], "dup": shapes_with_tables1[0].table.rows[6].cells[0]},
                "resp": {"orig": shapes_with_tables[0].table.rows[9].cells[1], "dup": shapes_with_tables1[0].table.rows[9].cells[1]},
                "due": {"orig": shapes_with_tables[0].table.rows[9].cells[3], "dup": shapes_with_tables1[0].table.rows[9].cells[3]}
            }

            # Update cell texts and formatting for overflowed cells
            for key, info in overflow_info.items():
                # Original slide: show first part
                cell_map[key]["orig"].text = info["first"]
                format_cell(cell_map[key]["orig"])
                # Duplicated slide: show remaining part
                cell_map[key]["dup"].text = info["remaining"]
                format_cell(cell_map[key]["dup"]) # '''
                
        # split_and_duplicate_slides(
        #     ppt, 
        #     new_slide, 
        #     shapes_with_tables, 
        #     shapes_with_textbox, 
        #     new_slide_idx
        # )

    # Deleting the original template slide
    ppt = delete_slide(ppt, 6)
    # Save the final presentation
    
    # Define the output folder inside the parent directory
    output_folder_path = Path(output_folder_path)
    output_folder = output_folder_path / "Internal Audit Report"
    # Create the folder if it doesn't exist
    output_folder.mkdir(parents=True, exist_ok=True)

    OP_DIR = output_folder
    if not os.path.exists(OP_DIR):
            os.makedirs(OP_DIR)
    output_path = r"%s/IFSL - IA Q4 - Collections  Recovery_Automated.pptx" % (OP_DIR)
    ppt.save(output_path)

    return output_path

if __name__ == "__main__":
    IA_Report_Utility_Main(r"C:\Users\satyambarnwal\OneDrive - KPMG\Desktop\Report Utility Misc files\IA Template 2 requirements\Input file for Report Automation.xlsx" , r"C:\Users\satyambarnwal\OneDrive - KPMG\Desktop\Report Utility Misc files\IA Template 2 requirements\Zensar_IA_Subcontracting Review_Q3 FY25-26 (Automation) - Template.pptx" , r"C:\Users\satyambarnwal\OneDrive - KPMG\Desktop\Report Utility Misc files\IA Template 2 requirements")


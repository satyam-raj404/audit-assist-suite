# inspect_placeholders_quick.py
import re
from pathlib import Path
from pptx import Presentation
import io
import re
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
import os
from datetime import datetime

from IA_Report_Template2.ppt_helpers import (_parse_color , format_text_in_table_cell , replace_rectangle_with_table_and_handle_overflow , duplicate_slide , 
                         split_text_by_max_lines , handle_cell_overflow , format_text_in_shape , check_overflow_text , fill_template_with_header_map_and_title ,
                         make_ordered_hashmap , map_columns , _find_table_for_mapping , remove_whitespace_cells , ann_replace_rectangle_with_table_and_handle_overflow ,
                         update_page_no ,
                         fill_templates_with_header_map_and_title_v2 ,
                         fill_detailed_observation_templates)

from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx.dml.color import RGBColor
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt

def _to_rgb_color(color):  # Aryan Addition
    """
    Convert a '#RRGGBB' hex string or (R, G, B) tuple into an RGBColor.
    """
    if isinstance(color, str):
        s = color.strip().lstrip('#')
        if len(s) != 6 or any(c not in "0123456789aAbBcCdDeEfF" for c in s):
            raise ValueError(f"Invalid hex color: {color!r}")
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
        return RGBColor(r, g, b)
    elif isinstance(color, (tuple, list)) and len(color) == 3:
        r, g, b = color
        if not all(isinstance(x, int) and 0 <= x <= 255 for x in (r, g, b)):
            raise ValueError(f"Invalid RGB tuple: {color!r}")
        return RGBColor(r, g, b)
    else:
        raise TypeError("color must be '#RRGGBB' or (R, G, B)")

def set_table_cell_format(
    cell,
    *,
    horizontal_align: str | None = "center",   # 'left' | 'center' | 'right' | 'justify' | 'distribute'
    vertical_align: str | None = "middle",     # 'top' | 'middle' | 'bottom'
    font_size_pt: float | None = 12,
    font_color = "#000000",                    # '#RRGGBB' or (R, G, B), None to skip
    bold: bool | None = True,                  # <<< set bold text (True/False/None)
    italic: bool | None = None
):
    """
    Apply alignment, font size, font color, and bold to a python-pptx table cell.

    Parameters
    ----------
    cell : pptx.table._Cell
    horizontal_align : str | None
        'left' | 'center' | 'right' | 'justify' | 'distribute' | None
    vertical_align : str | None
        'top' | 'middle' | 'bottom' | None
    font_size_pt : float | None
        Size in points; None leaves unchanged.
    font_color : str | tuple | None
        '#RRGGBB' or (R,G,B); None leaves unchanged.
    bold : bool | None
        True makes text bold, False removes bold, None leaves unchanged.
    italic : bool | None
        True/False/None similar to bold.
    """
    # --- vertical alignment (cell-level) ---
    if vertical_align:
        va_map = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        key = vertical_align.strip().lower()
        if key not in va_map:
            raise ValueError(f"vertical_align must be one of {list(va_map.keys())}")
        cell.vertical_anchor = va_map[key]

    tf = cell.text_frame  # ensure text frame

    # --- horizontal alignment (paragraph-level) ---
    if horizontal_align:
        ha_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
            "distribute": PP_ALIGN.DISTRIBUTE,
        }
        key = horizontal_align.strip().lower()
        if key not in ha_map:
            raise ValueError(f"horizontal_align must be one of {list(ha_map.keys())}")
        for para in tf.paragraphs:
            para.alignment = ha_map[key]

    # --- font attributes (run-level) ---
    color_rgb = _to_rgb_color(font_color) if font_color is not None else None
    size_pt = Pt(font_size_pt) if font_size_pt is not None else None

    for para in tf.paragraphs:
        # If paragraph has no runs, create one to apply formatting
        if not para.runs:
            para.add_run()
        for run in para.runs:
            font = run.font
            if size_pt is not None:
                font.size = size_pt
            if color_rgb is not None:
                font.color.rgb = color_rgb
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic

def set_table_cell_fill(cell, color, *, set_text_color=None):  #Aryan Addition
    """
    Set a python-pptx table cell's background fill to a solid color.

    Parameters
    ----------
    cell : pptx.table._Cell
        A table cell object (e.g., table.cell(row, col)).
    color : str | tuple[int, int, int]
        Fill color as '#RRGGBB' or (R, G, B).
    set_text_color : None | str | tuple[int, int, int]
        If provided, also set text (font) color in the cell.

    Notes
    -----
    - This uses a SOLID fill. (python-pptx doesn't support gradient fills.)
    - If the cell previously had a theme/gradient/pattern fill, we clear and
      apply a solid fill.
    """
    rgb = _to_rgb_color(color)

    # Ensure fill exists and is set to solid
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = rgb

    # Optionally set text color for all runs in the cell
    if set_text_color is not None:
        txt_rgb = _to_rgb_color(set_text_color)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = txt_rgb

# def delete_empty_rows(table): # Aryan Addition
#     # Remove rows where first cell is still "-"
#     tbl = table._tbl
#     # rows = list(table.rows)
#     i = 0
#     for row in table.rows:  # Exclude total row
#         if row.cells[0].text == "-":
#             tr = tbl.tr_lst[i]
#             row.cells[0].text = "Count"
#             # tbl.remove(tr)
#         i += 1

def delete_empty_rows(table, *, is_empty=None) -> int: # Aryan Addition
    """
    Remove all rows in a python-pptx table where *every* cell is empty.
    
    Parameters
    ----------
    table : pptx.table.Table
        The table object (e.g., shape.table) to process.
    is_empty : callable | None
        Optional predicate function taking `cell` and returning True if the cell
        should be considered empty. Defaults to `str(cell.text).strip()==""`.
    
    Returns
    -------
    int
        The number of rows removed.
    """
    # Default "empty" definition: no text or only whitespace
    def _default_is_empty(cell) -> bool:
        txt = cell.text if cell.text is not None else ""
        return txt.strip() == ""
    
    empty_pred = is_empty or _default_is_empty

    n_rows = len(table.rows)
    n_cols = len(table.columns)

    # Find rows that are empty across all columns
    empty_row_idxs = []
    for r in range(n_rows):
        if all(empty_pred(table.cell(r, c)) for c in range(n_cols)):
            empty_row_idxs.append(r)

    # Remove rows bottom-up to avoid index shifting
    removed = 0
    for r in reversed(empty_row_idxs):
        tr = table._tbl.tr_lst[r]   # underlying <a:tr> element
        table._tbl.remove(tr)
        removed += 1

    return removed

def replace_text_on_slide_simple(slide, old, new, exact=True):
    """
    Replace text on a slide.
    - slide: ppt.slides[index]
    - old: string to match
    - new: replacement string
    - exact: if True replace only when the whole shape text equals `old`;
             if False replace occurrences inside runs.
    """
    def _replace_in(sh):
        if getattr(sh, "has_text_frame", False):
            tf = sh.text_frame
            # quick exact-match path
            if exact:
                if (tf.text or "").strip() == old:
                    tf.clear()
                    tf.text = new
            else:
                for p in tf.paragraphs:
                    for run in p.runs:
                        if old in run.text:
                            run.text = run.text.replace(old, new)
        elif getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for s in sh.shapes:
                _replace_in(s)

    for shape in slide.shapes:
        _replace_in(shape)

from pptx.util import Pt
from pptx.dml.color import RGBColor
import pandas as pd

def write_observations_to_slide_table_v2_fixed(slide, df_obs: pd.DataFrame):
    """
    Fixed version of write_observations_to_slide_table_v2.

    - Uses explicit mapping:
        {"Observation": (2,0),
         "D":           (2,2),
         "IT":          (2,3),
         "OE":          (2,4),
         "O":           (2,5),
         "FR":          (2,6),
         "C":           (2,7)}
    - Each DataFrame row i is written to table row (base_row + i).
    - Does NOT change existing PPT cells when DF value is empty/NaN.
    - Does NOT duplicate slides or add rows; out-of-bounds target cells are skipped.
    - Applies Arial 10pt to inserted text.
    - Applies fill color for category values "H","M","L".
    - Uses "Issue heading" (case-insensitive) as the source column for Observation.
    - Returns number of DataFrame rows that had at least one cell written.
    """
    if df_obs is None or len(df_obs) == 0:
        return 0

    # find first table on slide
    table_shape = None
    for sh in slide.shapes:
        if getattr(sh, "has_table", False):
            table_shape = sh
            break
    if table_shape is None:
        raise ValueError("No table found on the provided slide.")

    tbl = table_shape.table

    # explicit mapping as requested
    mapping = {
        "Observation": (2, 0),
        "D":           (2, 2),
        "IT":          (2, 3),
        "OE":          (2, 4),
        "O":           (2, 5),
        "FR":          (2, 6),
        "C":           (2, 7),
    }

    # prepare case-insensitive column lookup for df_obs
    df = df_obs.reset_index(drop=True).copy()
    lower_to_actual = {c.lower(): c for c in df.columns}

    # Use "Issue heading" (case-insensitive) as the Observation source column
    obs_key_lower = "issue heading"
    obs_col_actual = None
    if obs_key_lower in lower_to_actual:
        obs_col_actual = lower_to_actual[obs_key_lower]
    elif "issueheading" in lower_to_actual:  # tolerate missing space
        obs_col_actual = lower_to_actual["issueheading"]
    else:
        # fallback: any column containing 'issue' and 'head' or 'heading' or 'observ'
        for c in df.columns:
            cl = c.lower()
            if "issue" in cl and ("head" in cl or "heading" in cl):
                obs_col_actual = c
                break
        if obs_col_actual is None:
            for c in df.columns:
                if "observ" in c.lower():
                    obs_col_actual = c
                    break

    # build final column name mapping (map keys -> actual df column names or None)
    final_col_map = {}
    for key in mapping.keys():
        if key == "Observation":
            final_col_map[key] = obs_col_actual
        else:
            # case-insensitive match for other keys
            final_col_map[key] = key if key in df.columns else lower_to_actual.get(key.lower())

    # color mapping for category values
    color_map = {
        "H": RGBColor(0xFF, 0x00, 0x00),
        "M": RGBColor(0xFF, 0xC0, 0x00),
        "L": RGBColor(0x92, 0xD0, 0x50),
    }

    rows_with_writes = 0

    for i in range(len(df)):
        df_row = df.iloc[i]
        wrote_any_cell = False

        for col_key, (base_row, base_col) in mapping.items():
            target_row = base_row + i
            target_col = base_col

            # bounds check
            if target_row < 0 or target_row >= len(tbl.rows):
                continue
            if target_col < 0 or target_col >= len(tbl.columns):
                continue

            # determine actual df column name
            src_col = final_col_map.get(col_key)
            if not src_col:
                # no matching column in df for this key; skip
                continue

            # fetch value from df safely
            try:
                val = df_row[src_col]
            except Exception:
                val = ""

            # skip empty/NaN -> leave existing PPT cell unchanged
            if pd.isna(val) or (isinstance(val, str) and val.strip() == ""):
                continue

            text_to_write = str(val)

            # write into PPT cell with robust text-frame handling
            try:
                cell = tbl.rows[target_row].cells[target_col]
            except Exception:
                continue

            try:
                # If text_frame exists, clear and set text via paragraph to preserve formatting control
                tf = getattr(cell, "text_frame", None)
                if tf is not None:
                    try:
                        tf.clear()
                    except Exception:
                        # some pptx versions require clearing paragraphs manually
                        try:
                            while tf.paragraphs:
                                tf.paragraphs.pop()
                        except Exception:
                            pass
                    # ensure at least one paragraph exists
                    if not tf.paragraphs:
                        p = tf.add_paragraph()
                    p = tf.paragraphs[0]
                    # set text on paragraph (safer than manipulating runs directly)
                    p.text = text_to_write
                    # ensure a run exists and set font on the first run
                    try:
                        run = p.runs[0] if p.runs else p.add_run()
                        run.text = text_to_write
                        run.font.name = "Arial"
                        run.font.size = Pt(10)
                        run.font.bold = False
                    except Exception:
                        # fallback: set paragraph text only
                        p.text = text_to_write
                else:
                    # fallback: direct assignment
                    cell.text = text_to_write
                wrote_any_cell = True
            except Exception:
                # fallback attempt: direct assignment
                try:
                    cell.text = text_to_write
                    wrote_any_cell = True
                except Exception:
                    pass

            # apply color only for category columns (D, IT, OE, O, FR, C)
            if col_key in {"D", "IT", "OE", "O", "FR", "C"}:
                v_norm = text_to_write.strip().upper()
                if v_norm in color_map:
                    try:
                        fill = cell.fill
                        fill.solid()
                        fill.fore_color.rgb = color_map[v_norm]
                    except Exception:
                        pass

        if wrote_any_cell:
            rows_with_writes += 1
    
    row_idx = 2
    for idx, value in df["RR"].items():
        cell = tbl.cell(row_idx, 1)
        if value == 'H':
            cell.text = 'High'
            set_table_cell_fill(cell, "#FF0000")            
            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=12, font_color="#FFFFFF", bold=True)
        
        elif value == 'L':
            cell.text = 'Low'
            set_table_cell_fill(cell, '#92D050')
            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=12, font_color="#000000", bold=True)
        
        elif value == 'M':
            cell.text = 'Medium'
            set_table_cell_fill(cell, '#FFC000')
            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=12, font_color="#000000", bold=True)

        else:
            continue
        
        row_idx += 1
        
    delete_empty_rows(tbl)

    return rows_with_writes
#--------------------------------------------------------------------------------------------------
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

def apply_table_style(tbl):
    """Style a pptx table: header fill, Arial 10.5, left align, spacing, line spacing, gray borders."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls

    header_color = RGBColor(0, 51, 141)
    white = RGBColor(255, 255, 255)
    font_name = "Arial"
    font_size = Pt(10.5)
    border_val = "{:02X}{:02X}{:02X}".format(128, 128, 128)
    border_w = "12700"  # ~0.5pt

    # set column widths and row heights if desired externally; here we only style cells
    rows = len(tbl.rows)
    cols = len(tbl.columns)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.rows[r].cells[c]
            text = (cell.text or "")
            cell.text = text
            # vertical center
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # paragraph formatting (use direct paragraph attributes; avoid paragraph_format access)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            p.line_spacing = 1.2
            # ensure a run exists and set font
            run = p.runs[0] if p.runs else p.add_run()
            run.text = p.text
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = (r == 0)
            # fills
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = white
            # borders via XML
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            xml = f"""
            <a:tcPr {nsdecls('a')}>
              <a:lnL w="{border_w}"><a:solidFill><a:srgbClr val="{border_val}"/></a:solidFill></a:lnL>
              <a:lnR w="{border_w}"><a:solidFill><a:srgbClr val="{border_val}"/></a:solidFill></a:lnR>
              <a:lnT w="{border_w}"><a:solidFill><a:srgbClr val="{border_val}"/></a:solidFill></a:lnT>
              <a:lnB w="{border_w}"><a:solidFill><a:srgbClr val="{border_val}"/></a:solidFill></a:lnB>
            </a:tcPr>
            """
            new_tcPr = parse_xml(xml)
            for child in new_tcPr:
                tcPr.append(child)

def insert_key_stats_table(ppt, df, template_slide, rect_shape_idx=None):
    """
    Create one duplicate of template_slide, replace its main rectangle with a table built from df,
    and apply apply_table_style(tbl). Returns the new slide object.
    Parameters:
      - ppt: Presentation object
      - df: pandas DataFrame for a single "Key stats" sheet (raw or cleaned)
      - template_slide: Slide object to duplicate (e.g., ppt.slides[9])
      - rect_shape_idx: optional index of rectangle to replace on the duplicated slide
    """
    # minimal cleaning: strip strings, drop empty rows/cols
    import pandas as pd
    d = df.copy()
    d.columns = [str(c).strip() for c in d.columns]
    def _clean(v):
        if pd.isna(v): return ""
        s = str(v).strip()
        if s in ("-", "nan", "None"): return ""
        return s
    d = d.applymap(_clean)
    d = d.loc[d.apply(lambda r: any(r.astype(bool)), axis=1)].reset_index(drop=True)
    d = d[[c for c in d.columns if any(d[c].astype(bool))]]

    # duplicate slide by recreating shapes
    def _duplicate_after(pres, src):
        new = pres.slides.add_slide(src.slide_layout)
        for sh in src.shapes:
            if getattr(sh, "has_text_frame", False):
                left, top, w, h = sh.left, sh.top, sh.width, sh.height
                tx = new.shapes.add_textbox(left, top, w, h)
                dst = tx.text_frame
                dst.clear()
                for p in sh.text_frame.paragraphs:
                    np = dst.add_paragraph()
                    np.level = p.level
                    for run in p.runs:
                        nr = np.add_run()
                        nr.text = run.text
                        if run.font is not None:
                            if run.font.name: nr.font.name = run.font.name
                            if run.font.size: nr.font.size = run.font.size
                            nr.font.bold = bool(run.font.bold)
                            if getattr(run.font, "color", None) and getattr(run.font.color, "rgb", None):
                                nr.font.color.rgb = run.font.color.rgb
            elif getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
                left, top, w, h = sh.left, sh.top, sh.width, sh.height
                tbl = sh.table
                rows, cols = len(tbl.rows), len(tbl.columns)
                new_tbl_sh = new.shapes.add_table(rows=rows, cols=cols, left=left, top=top, width=w, height=h)
                new_tbl = new_tbl_sh.table
                for r in range(rows):
                    for c in range(cols):
                        src_cell = tbl.rows[r].cells[c]
                        dst_cell = new_tbl.rows[r].cells[c]
                        dst_cell.text = src_cell.text_frame.text if getattr(src_cell, "text_frame", None) else ""
            elif getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                img = sh.image
                new.shapes.add_picture(img.blob, sh.left, sh.top, width=sh.width, height=sh.height)
        sldIdLst = pres.slides._sldIdLst
        new_id = sldIdLst[-1]
        src_idx = list(pres.slides).index(src)
        sldIdLst.remove(new_id)
        sldIdLst.insert(src_idx + 1, new_id)
        return pres.slides[src_idx + 1]

    def _find_rectangle(slide):
        best = None
        best_area = 0
        for sh in slide.shapes:
            w = getattr(sh, "width", 0)
            h = getattr(sh, "height", 0)
            area = int(w) * int(h)
            if getattr(sh, "has_text_frame", False) or getattr(sh, "shape_type", None) in (1, 17, 19):
                if area > best_area:
                    best = sh
                    best_area = area
        return best

    new_slide = _duplicate_after(ppt, template_slide)
    rect = new_slide.shapes[rect_shape_idx] if rect_shape_idx is not None else _find_rectangle(new_slide)
    left, top, width, height = rect.left, rect.top, rect.width, rect.height
    new_slide.shapes._spTree.remove(rect._element)

    # determine header and data
    if d.shape[0] > 0 and all(str(x).strip() != "" for x in d.iloc[0].tolist()):
        # header = [str(x) for x in d.iloc[0].tolist()]
        header = list(df.columns)
        data = d.iloc[1:].reset_index(drop=True)
    else:
        header = [str(c) for c in d.columns]
        data = d.reset_index(drop=True)

    rows_needed = 1 + len(data)
    cols = max(1, len(header))
    tbl_sh = new_slide.shapes.add_table(rows=rows_needed, cols=cols, left=left, top=top, width=width, height=height)
    tbl = tbl_sh.table

    # populate header and data
    for c, h in enumerate(header[:cols]):
        tbl.rows[0].cells[c].text = "" if h is None else str(h)
    for r_idx in range(len(data)):
        row = data.iloc[r_idx].tolist()
        for c_idx in range(min(len(row), cols)):
            tbl.rows[1 + r_idx].cells[c_idx].text = "" if row[c_idx] is None else str(row[c_idx])

    # # per‑cell sizes (points)
    # col_pt = Pt(20)       # standard column width
    # row_pt = Pt(35)      # standard row height
    # # determine header and data first 
    # rows_needed = 1 + len(data) # header + data rows 
    # cols = max(1, len(header))
    # # enforce exact column widths and row heights
    # for c in range(cols):
    #     tbl.columns[c].width = col_pt
    # for r in range(rows_needed):
    #     tbl.rows[r].height = row_pt

    apply_table_style(tbl)
    return new_slide
###################################################################################################

###################################################################################################
#--------------------------------------------------------------------------------------------------
def IA_Report_Utility_MainT2(tracker_path, ppt_path, output_folder_path_og):

    output_folder_path = output_folder_path_og + '\IA-Zensar'

    # 1) Ensure the folder exists
    os.makedirs(output_folder_path, exist_ok=True)

    # 2) Build a safe output file path
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = os.path.join(output_folder_path, f"IA-Zensar-Automated_{timestamp}.pptx")


    # Importing the Issue Tracker as a dictionary
    raw_df = pd.read_excel(tracker_path , sheet_name=None , engine="openpyxl")
    
    # Creating seperate dataframes from the dictionary by taking the sheet name as keys (Not used as much, used for testing purposes)
    df_obs = raw_df["Observation listing"]
    df_key1 = raw_df["Key Stats"]
    df_key2 = raw_df["Key Stats-1"]
    df_key3 = raw_df["Key Stats-2"]
    df_key4 = raw_df["Key Stats-3"]
    df_key5 = raw_df["Key Stats-4"]
    df_keyan1 = raw_df["Annexure I.a"]

    # header = map_columns(df_obs)

    # Making hashmaps for the titles, headers for tables, etc.
    Header_map = make_ordered_hashmap(df_obs)

    # Importing the ppt as a python object
    ppt = Presentation(r"%s" % (ppt_path))

    # Defining separate variables for template slides
    template5 = ppt.slides[6] # List of Observations
    template7 = ppt.slides[8] # Executive Summary
    template11 = ppt.slides[12] # Detailed Observation 1
    template12 = ppt.slides[13] # Detailed Observation 2
    template9 = ppt.slides[10] # Stats
    template14 = ppt.slides[15] # Annexure

    # Inserting the Key-Stats (Overview and Key Infromation) slides (Slide 10)
    insert_key_stats_table(ppt, remove_whitespace_cells(df_key4), template9)

    # Creating indexes for where the cells are in the table (x,y) where x is rows and y is columns 
    HEADER_CELL_MAP7 = {"Observations": (1, 0), "Recommendation": (1, 1)}
    HEADER_CELL_MAP11 = {"Observations": (1, 0), "Root cause": (3, 0)}
    HEADER_CELL_MAP12 = {"Potential risk": (1, 0), "Recommendation": (3, 0) , "Management Action Plan": (5, 0)}
    
    # Filling the Executive summary slide
    c = fill_template_with_header_map_and_title(raw_df, ppt, template_slide = template7, header_cell_map=HEADER_CELL_MAP7, header_map=Header_map, title_shape_idx=0 , end_char="a")

    # Function to fill risk mapping
    write_observations_to_slide_table_v2_fixed(template5 , df_obs) #for risk mapping

    # To fill the Key Stats
    # collect keys from raw_df that start with "Key Stats" (case-sensitive)
    key_list = [k for k in raw_df.keys() if str(k).startswith("Key Stats")]
    # optional: sort so "Key Stats", "Key Stats-1", "Key Stats-2", ... are in natural order
    def _sort_key(k):
        parts = k.split("-", 1)
        if len(parts) == 1:
            return (0,)
        try:
            return (1, int(parts[1]))
        except Exception:
            return (1, parts[1])
    key_list.sort(key=_sort_key)

    for key in key_list:
        df = remove_whitespace_cells(raw_df[key])
        try:
            insert_key_stats_table(ppt, df, template9)
        except Exception as e:
            print(f"Failed for {key}: {e}")

        ankey_list = [k for k in raw_df.keys() if str(k).startswith("Annexure")]
    def _sort_key(k):
        parts = k.split("-", 1)
        if len(parts) == 1:
            return (0,)
        try:
            return (1, int(parts[1]))
        except Exception:
            return (1, parts[1])
    ankey_list.sort(key=_sort_key)

    for key in ankey_list:
        df1 = remove_whitespace_cells(raw_df[key])
        try:
            insert_key_stats_table(ppt, df1, template14)
        except Exception as e:
            print(f"Failed for {key}: {e}")
    
    # Filling detailed observation slides
    fill_detailed_observation_templates(raw_df, ppt, template11, template12, list(ppt.slides).index(template12)+1)

    # fill_templates_with_header_map_and_title_v2(
    #     raw_df=raw_df, ppt=ppt, template_slides=[template11, template12],   # ← pass both templates here
    #     header_cell_map=HEADER_CELL_MAP11, header_map=Header_map, title_shape_idx=0,
    #     title_template="{{Sl No.}}. {{Issue heading}}", title_font_pt=24, title_bold=True, cell_font_pt=10, cell_bold=False, end_char="x"   # optional; returns ending slide number (of the last processed template)
    # )

    # d = fill_template_with_header_map_and_title(raw_df, ppt, template_slide = template11, header_cell_map=HEADER_CELL_MAP11, header_map=Header_map, title_shape_idx=0 , end_char="b")

    # # Filling Detailed observation slide 2
    # e = fill_template_with_header_map_and_title(raw_df, ppt, template_slide = template12, header_cell_map=HEADER_CELL_MAP12, header_map=Header_map, title_shape_idx=0 , end_char="c") 

    # Updating page numbers
    update_page_no(ppt, 1, "Audit Scope", "a")
    update_page_no(ppt, 1, "List of Observations", "b")
    update_page_no(ppt, 1, "Executive Summary", "c")
    update_page_no(ppt, 1, "Overview and Key Information", "d")
    update_page_no(ppt, 1, "Detailed Observations", "e")
    update_page_no(ppt, 1, "Annexures", "f")
            

    ppt.save(output_file)


# if __name__ == "__main__" :
#     tracker_path = r"IA_Report_Template2\Input\Input file for Report Automation.xlsx"
#     ppt_path = r"IA_Report_Template2\Input\Zensar_IA_Subcontracting Review_Q3 FY25-26 (Automation) - Template.pptx"
#     output_folder_path = r"IA_Report_Template2\Output\IA_Report_Automation.pptx"

#     IA_Report_Utility_MainT2(tracker_path, ppt_path, output_folder_path)



#ppt.slides[7].shapes[1].table.rows[1].cells[0]--- Locate using theses indexes
#annextures adding in the ppt using this code snippet-------------------------------------------------
#df = raw_df.get("Annexure I.a") 
# replace_rectangle_with_table_and_handle_overflow( ppt, template_slide_idx=14, sheet_df=df, rect_shape_idx=None, max_rows_per_slide=7, max_lines_per_cell=14 )
#-----------------------------------------------------------------------------------------------------
 
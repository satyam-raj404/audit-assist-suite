from pptx.util import Inches
import math
from .Cell_text_insertion import duplicate_slide
from typing import Tuple, Optional
def _split_cell_text_by_lines(val, max_lines_per_cell: int) -> Tuple[bool, str, Optional[str]]:
    """
    Split a cell value into a first part that fits within `max_lines_per_cell`
    and a remaining part for overflow.

    Returns:
      (had_overflow: bool, first_part: str, remaining: Optional[str])

    Behavior:
    - Treats None and NaN as empty string (no overflow).
    - Converts non-string values to str().
    - Splits on newline characters; does NOT perform word-wrapping.
    - If max_lines_per_cell < 1, everything is treated as overflow.
    """
    if max_lines_per_cell is None:
        max_lines_per_cell = 0

    # normalize empty / NaN
    if val is None:
        return False, "", None
    try:
        if isinstance(val, float) and math.isnan(val):
            return False, "", None
    except Exception:
        pass

    s = str(val)

    # guard for non-positive max_lines_per_cell
    if max_lines_per_cell < 1:
        return True, "", s if s != "" else None

    lines = s.splitlines() or [s]

    if len(lines) <= max_lines_per_cell:
        return False, s, None

    first_part = "\n".join(lines[:max_lines_per_cell])
    remaining = "\n".join(lines[max_lines_per_cell:])

    return True, first_part, remaining

def _safe_remove_shape(slide, shape):
    """Remove a shape from slide safely (uses element removal)."""
    try:
        slide.shapes._spTree.remove(shape._element)
    except Exception:
        # fallback: try parent removal
        try:
            shape._element.getparent().remove(shape._element)
        except Exception:
            pass

# def _find_rectangle_shape(slide):
#     """
#     Heuristic: return first shape that has no table and is not the title and
#     has a non-zero width/height. If multiple, pick the largest by area.
#     """
#     candidates = []
#     for sh in slide.shapes:
#         if getattr(sh, "has_table", False):
#             continue
#         # skip title placeholder
#         if sh == getattr(slide.shapes, "title", None):
#             continue
#         # require numeric geometry
#         try:
#             w = int(sh.width)
#             h = int(sh.height)
#         except Exception:
#             continue
#         if w <= 0 or h <= 0:
#             continue
#         candidates.append((w * h, sh))
#     if not candidates:
#         return None
#     # pick largest area
#     candidates.sort(key=lambda x: x[0], reverse=True)
#     return candidates[0][1]


# #for key stats
# def replace_rectangle_with_table_and_handle_overflow(
#         ppt,
#         template_slide_idx,
#         sheet_df,
#         rect_shape_idx=None,
#         max_rows_per_slide=20,
#         max_lines_per_cell=14):
#     """
#     Replace rectangle on template slide with a table built from sheet_df.
#     If data overflows the table capacity or per-cell line limit, create duplicates
#     using duplicate_slide(ppt, template_slide) and place overflow rows there.
#     - ppt: Presentation object
#     - template_slide_idx: index of the template slide (e.g., 14)
#     - sheet_df: pandas DataFrame for the annexure (e.g., raw_df['Annexure I.a'])
#     - rect_shape_idx: optional shape index to replace; if None, heuristic finds the rectangle
#     - max_rows_per_slide: total rows in PPT table including header (default 20)
#     - max_lines_per_cell: per-cell line threshold (default 14)
#     """
#     # basic validation
#     if sheet_df is None or len(sheet_df) == 0:
#         print("No data in sheet_df; nothing to write.")
#         return

#     template_slide = ppt.slides[template_slide_idx]

#     # locate rectangle shape
#     rect_shape = None
#     if rect_shape_idx is not None:
#         try:
#             rect_shape = template_slide.shapes[rect_shape_idx]
#         except Exception:
#             rect_shape = None
#     if rect_shape is None:
#         rect_shape = _find_rectangle_shape(template_slide)
#     if rect_shape is None:
#         raise ValueError("Could not find a rectangle shape to replace on template slide")

#     # capture geometry and remove the rectangle
#     left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height
#     _safe_remove_shape(template_slide, rect_shape)

#     # helper to add a table in the same bounds
#     def _add_table_on_slide(slide, df, rows_total):
#         cols = max(1, len(df.columns))
#         tbl_sh = slide.shapes.add_table(rows=rows_total, cols=cols,
#                                        left=left, top=top, width=width, height=height)
#         tbl = tbl_sh.table
#         # header row
#         for c, col in enumerate(df.columns):
#             if c >= len(tbl.columns):
#                 break
#             tbl.rows[0].cells[c].text = str(col)
#         # set proportional column widths (simple equal split if many columns)
#         try:
#             col_width = width // max(1, cols)
#             for c in range(min(cols, len(tbl.columns))):
#                 tbl.columns[c].width = col_width
#         except Exception:
#             pass
#         return tbl_sh

#     # prepare remaining_df (copy)
#     remaining_df = sheet_df.reset_index(drop=True).copy()

#     # loop to create first table on template and duplicates as needed
#     last_insert_id = None
#     first_iteration = True
#     while len(remaining_df) > 0:
#         # on first iteration, add table to template_slide; on subsequent, duplicate template and use that
#         if first_iteration:
#             target_slide = template_slide
#             first_iteration = False
#             # create table on template slide
#             rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#             table_shape = _add_table_on_slide(target_slide, remaining_df, rows_to_create)
#         else:
#             # duplicate the template slide (which now contains the table)
#             new_slide = duplicate_slide(ppt, template_slide)
#             # set title to sheet name if title exists
#             if new_slide.shapes.title is not None:
#                 new_slide.shapes.title.text = str(getattr(sheet_df, "name", "") or "")
#             # newly appended slide is last; move it immediately after last_insert_id or template
#             sldIdLst = ppt.slides._sldIdLst
#             new_id = sldIdLst[-1]
#             anchor_id = last_insert_id if last_insert_id is not None else sldIdLst[template_slide_idx]
#             sldIdLst.remove(new_id)
#             sldIdLst.insert(sldIdLst.index(anchor_id) + 1, new_id)
#             last_insert_id = new_id
#             # find the table on the duplicated slide (first table)
#             target_slide = ppt.slides[sldIdLst.index(new_id)]
#             table_shape = None
#             for sh in target_slide.shapes:
#                 if getattr(sh, "has_table", False):
#                     table_shape = sh
#                     break
#             if table_shape is None:
#                 # if duplication didn't copy the table for some reason, add one in same bounds
#                 rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#                 table_shape = _add_table_on_slide(target_slide, remaining_df, rows_to_create)

#         tbl = table_shape.table
#         capacity = len(tbl.rows) - 1  # data rows capacity (row 0 is header)
#         # prepare chunk to write
#         chunk = remaining_df.iloc[:capacity].reset_index(drop=True)
#         rows_written = 0
#         overflow_map = {}  # (row_in_chunk, col_idx) -> remaining_text

#         for r_idx, (_, row) in enumerate(chunk.iterrows()):
#             if r_idx >= capacity:
#                 break
#             for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
#                 cell = tbl.rows[1 + r_idx].cells[c_idx]
#                 overflow, first_part, remaining = _split_cell_text_by_lines(val, max_lines_per_cell)
#                 cell.text = first_part
#                 if remaining:
#                     overflow_map[(r_idx, c_idx)] = remaining
#             rows_written += 1

#         # drop written rows from remaining_df
#         remaining_df = remaining_df.iloc[rows_written:].reset_index(drop=True)

#         # if there are per-cell overflows, create a single-row DataFrame to prepend to remaining_df
#         if overflow_map:
#             cols = list(sheet_df.columns)
#             overflow_row = {c: "" for c in cols}
#             for (r_in_chunk, c_idx), rem_text in overflow_map.items():
#                 if c_idx < len(cols):
#                     overflow_row[cols[c_idx]] = rem_text
#             # prepend overflow row so it is written first on next slide
#             remaining_df = pd.concat([pd.DataFrame([overflow_row]), remaining_df], ignore_index=True)

#         # if there are still rows remaining, ensure next iteration will insert after this slide
#         if len(remaining_df) > 0 and last_insert_id is None:
#             # if we haven't set last_insert_id yet (no duplicates created), set it to template id
#             sldIdLst = ppt.slides._sldIdLst
#             last_insert_id = sldIdLst[template_slide_idx]

#     print("Completed writing annexure table; duplicates created as needed.")

# def replace_rectangle_with_table_and_handle_overflow(
#         ppt,
#         template_slide,               # now expects a pptx.slide.Slide object
#         sheet_df,
#         rect_shape_idx=None,
#         max_rows_per_slide=20,
#         max_lines_per_cell=14):
#     """
#     Replace rectangle on template_slide with a table built from sheet_df.
#     - template_slide: pptx.slide.Slide object (not an index)
#     - sheet_df: pandas DataFrame
#     - rect_shape_idx: optional shape index to replace; if None, heuristic finds the rectangle
#     - max_rows_per_slide: total rows in PPT table including header (default 20)
#     - max_lines_per_cell: per-cell line threshold (default 14)
#     """
#     # basic validation
#     if sheet_df is None or len(sheet_df) == 0:
#         print("No data in sheet_df; nothing to write.")
#         return

#     # normalize template slide index for sldId list operations
#     template_index = list(ppt.slides).index(template_slide)

#     # locate rectangle shape
#     rect_shape = None
#     if rect_shape_idx is not None:
#         try:
#             rect_shape = template_slide.shapes[rect_shape_idx]
#         except Exception:
#             rect_shape = None
#     if rect_shape is None:
#         rect_shape = _find_rectangle_shape(template_slide)
#     if rect_shape is None:
#         raise ValueError("Could not find a rectangle shape to replace on template slide")

#     # capture geometry and remove the rectangle
#     left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height
#     _safe_remove_shape(template_slide, rect_shape)

#     # helper to add a table in the same bounds
#     def _add_table_on_slide(slide, df, rows_total):
#         cols = max(1, len(df.columns))
#         tbl_sh = slide.shapes.add_table(rows=rows_total, cols=cols,
#                                        left=left, top=top, width=width, height=height)
#         tbl = tbl_sh.table
#         # header row
#         for c, col in enumerate(df.columns):
#             if c >= len(tbl.columns):
#                 break
#             tbl.rows[0].cells[c].text = str(col)
#         # set proportional column widths (simple equal split)
#         col_width = width // max(1, cols)
#         for c in range(min(cols, len(tbl.columns))):
#             tbl.columns[c].width = col_width
#         return tbl_sh

#     # prepare remaining_df (copy)
#     remaining_df = sheet_df.reset_index(drop=True).copy()

#     # loop to create first table on template and duplicates as needed
#     last_insert_id = None
#     first_iteration = True
#     while len(remaining_df) > 0:
#         # on first iteration, add table to template_slide; on subsequent, duplicate template and use that
#         if first_iteration:
#             target_slide = template_slide
#             first_iteration = False
#             rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#             table_shape = _add_table_on_slide(target_slide, remaining_df, rows_to_create)
#         else:
#             # duplicate the template slide (pass Slide object)
#             new_slide = duplicate_slide(ppt, template_slide)
#             # set title to sheet name if title exists
#             if new_slide.shapes.title is not None:
#                 new_slide.shapes.title.text = str(getattr(sheet_df, "name", "") or "")
#             # newly appended slide is last; move it immediately after last_insert_id or template
#             sldIdLst = ppt.slides._sldIdLst
#             new_id = sldIdLst[-1]
#             anchor_id = last_insert_id if last_insert_id is not None else sldIdLst[template_index]
#             sldIdLst.remove(new_id)
#             sldIdLst.insert(sldIdLst.index(anchor_id) + 1, new_id)
#             last_insert_id = new_id
#             # find the table on the duplicated slide (first table)
#             new_slide_index = sldIdLst.index(new_id)
#             target_slide = ppt.slides[new_slide_index]
#             table_shape = None
#             for sh in target_slide.shapes:
#                 if getattr(sh, "has_table", False):
#                     table_shape = sh
#                     break
#             if table_shape is None:
#                 rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#                 table_shape = _add_table_on_slide(target_slide, remaining_df, rows_to_create)

#         tbl = table_shape.table
#         capacity = len(tbl.rows) - 1  # data rows capacity (row 0 is header)
#         # prepare chunk to write
#         chunk = remaining_df.iloc[:capacity].reset_index(drop=True)
#         rows_written = 0
#         overflow_map = {}  # (row_in_chunk, col_idx) -> remaining_text

#         for r_idx, (_, row) in enumerate(chunk.iterrows()):
#             if r_idx >= capacity:
#                 break
#             for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
#                 cell = tbl.rows[1 + r_idx].cells[c_idx]
#                 overflow, first_part, remaining = _split_cell_text_by_lines(val, max_lines_per_cell)
#                 cell.text = first_part or ""
#                 if remaining:
#                     overflow_map[(r_idx, c_idx)] = remaining
#             rows_written += 1

#         # drop written rows from remaining_df
#         remaining_df = remaining_df.iloc[rows_written:].reset_index(drop=True)

#         # if there are per-cell overflows, create a single-row DataFrame to prepend to remaining_df
#         if overflow_map:
#             cols = list(sheet_df.columns)
#             overflow_row = {c: "" for c in cols}
#             for (_, c_idx), rem_text in overflow_map.items():
#                 if c_idx < len(cols):
#                     overflow_row[cols[c_idx]] = rem_text
#             # prepend overflow row so it is written first on next slide
#             remaining_df = pd.concat([pd.DataFrame([overflow_row]), remaining_df], ignore_index=True)

#         # if there are still rows remaining, ensure next iteration will insert after this slide
#         if len(remaining_df) > 0 and last_insert_id is None:
#             sldIdLst = ppt.slides._sldIdLst
#             last_insert_id = sldIdLst[template_index]

#     print("Completed writing annexure table; duplicates created as needed.")

from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd

def _find_rectangle_shape(slide: Slide):
    """
    Heuristic: find the largest rectangle-like shape on the slide.
    - Considers top-level shapes and children inside groups.
    - Prefers AUTO_SHAPE and TEXT_BOX types, ignores shapes that contain tables.
    - Returns the Shape object with the largest area or None if none found.
    """
    candidates = []

    def consider(sh):
        if getattr(sh, "has_table", False):
            return
        # skip title placeholder
        if sh == getattr(slide.shapes, "title", None):
            return
        try:
            w = int(sh.width)
            h = int(sh.height)
        except Exception:
            return
        if w <= 0 or h <= 0:
            return
        # prefer auto-shapes and text boxes slightly by adding a small boost
        boost = 0
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            boost += 1000000
        if sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            boost += 500000
        area = w * h + boost
        candidates.append((area, sh))

    for sh in slide.shapes:
        consider(sh)
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            # iterate group children
            try:
                for child in sh.shapes:
                    consider(child)
            except Exception:
                pass

    if not candidates:
        return None
    candidates.sort(key=lambda x: x[0], reverse=True)
    return candidates[0][1]

#ppt.slides._sldIdLst.remove(ppt.slides._sldIdLst[0])

# def replace_rectangle_with_table_and_handle_overflow(
#         ppt,
#         template_slide,               # expects a pptx.slide.Slide object
#         sheet_df: pd.DataFrame,
#         rect_shape_idx=None,
#         max_rows_per_slide=20,
#         max_lines_per_cell=14):
#     """
#     Replace the largest rectangle on template_slide with a table built from sheet_df.
#     - template_slide: Slide object (not index)
#     - If rect_shape_idx is provided, that shape is used; otherwise _find_rectangle_shape is used.
#     - Keeps original overflow/duplication logic.
#     """
#     if sheet_df is None or len(sheet_df) == 0:
#         print("No data in sheet_df; nothing to write.")
#         return

#     # get template slide index for sldId list operations
#     template_index = list(ppt.slides).index(template_slide)

#     # locate rectangle shape
#     rect_shape = None
#     if rect_shape_idx is not None:
#         try:
#             rect_shape = template_slide.shapes[rect_shape_idx]
#         except Exception:
#             rect_shape = None

#     if rect_shape is None:
#         rect_shape = _find_rectangle_shape(template_slide)

#     if rect_shape is None:
#         # helpful debug message listing shapes on the slide
#         shapes_info = []
#         for i, sh in enumerate(template_slide.shapes):
#             shapes_info.append(f"#{i} name={getattr(sh,'name',None)} type={getattr(sh,'shape_type',None)} w={getattr(sh,'width',None)} h={getattr(sh,'height',None)}")
#         raise ValueError("Could not find a rectangle shape to replace on template slide. Slide shapes:\n" + "\n".join(shapes_info))

#     # capture geometry and remove the rectangle
#     left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height
#     _safe_remove_shape(template_slide, rect_shape)

#     def _add_table_on_slide(slide, df, rows_total):
#         cols = max(1, len(df.columns))
#         tbl_sh = slide.shapes.add_table(rows=rows_total, cols=cols,
#                                        left=left, top=top, width=width, height=height)
#         tbl = tbl_sh.table
#         # header row
#         for c, col in enumerate(df.columns):
#             if c >= len(tbl.columns):
#                 break
#             tbl.rows[0].cells[c].text = str(col)
#         # equal column widths
#         col_width = width // max(1, cols)
#         for c in range(min(cols, len(tbl.columns))):
#             tbl.columns[c].width = col_width
#         return tbl_sh

#     remaining_df = sheet_df.reset_index(drop=True).copy()

#     last_insert_id = None
#     first_iteration = True

#     while len(remaining_df) > 0:
#         if first_iteration:
#             target_slide = template_slide
#             first_iteration = False
#             rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#             table_shape = _add_table_on_slide(target_slide, remaining_df, rows_to_create)
#         else:
#             new_slide = duplicate_slide(ppt, template_slide)
#             if new_slide.shapes.title is not None:
#                 new_slide.shapes.title.text = str(getattr(sheet_df, "name", "") or "")
#             sldIdLst = ppt.slides._sldIdLst
#             new_id = sldIdLst[-1]
#             anchor_id = last_insert_id if last_insert_id is not None else sldIdLst[template_index]
#             sldIdLst.remove(new_id)
#             sldIdLst.insert(sldIdLst.index(anchor_id) + 1, new_id)
#             last_insert_id = new_id
#             new_slide_index = sldIdLst.index(new_id)
#             target_slide = ppt.slides[new_slide_index]
#             table_shape = None
#             for sh in target_slide.shapes:
#                 if getattr(sh, "has_table", False):
#                     table_shape = sh
#                     break
#             if table_shape is None:
#                 rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#                 table_shape = _add_table_on_slide(target_slide, remaining_df, rows_to_create)

#         tbl = table_shape.table
#         capacity = len(tbl.rows) - 1
#         chunk = remaining_df.iloc[:capacity].reset_index(drop=True)
#         rows_written = 0
#         overflow_map = {}

#         for r_idx, (_, row) in enumerate(chunk.iterrows()):
#             if r_idx >= capacity:
#                 break
#             for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
#                 cell = tbl.rows[1 + r_idx].cells[c_idx]
#                 overflow, first_part, remaining = _split_cell_text_by_lines(val, max_lines_per_cell)
#                 cell.text = first_part or ""
#                 if remaining:
#                     overflow_map[(r_idx, c_idx)] = remaining
#             rows_written += 1

#         remaining_df = remaining_df.iloc[rows_written:].reset_index(drop=True)

#         if overflow_map:
#             cols = list(sheet_df.columns)
#             overflow_row = {c: "" for c in cols}
#             for (_, c_idx), rem_text in overflow_map.items():
#                 if c_idx < len(cols):
#                     overflow_row[cols[c_idx]] = rem_text
#             remaining_df = pd.concat([pd.DataFrame([overflow_row]), remaining_df], ignore_index=True)

#         if len(remaining_df) > 0 and last_insert_id is None:
#             sldIdLst = ppt.slides._sldIdLst
#             last_insert_id = sldIdLst[template_index]

#     print("Completed writing annexure table; duplicates created as needed.")

# def replace_rectangle_with_table_and_handle_overflow(
#         ppt,
#         template_slide,               # expects a pptx.slide.Slide object
#         sheet_df: pd.DataFrame,
#         rect_shape_idx=None,
#         max_rows_per_slide=20,
#         max_lines_per_cell=14):
#     """
#     Replace the largest rectangle on template_slide with a table built from sheet_df.
#     This variant does NOT duplicate slides. It writes only the minimum amount of data
#     that fits on a single slide and reduces the table row height so more rows fit visually.
#     """
#     if sheet_df is None or len(sheet_df) == 0:
#         print("No data in sheet_df; nothing to write.")
#         return

#     # get template slide index for potential sldId list operations (kept for compatibility)
#     template_index = list(ppt.slides).index(template_slide)

#     # locate rectangle shape
#     rect_shape = None
#     if rect_shape_idx is not None:
#         try:
#             rect_shape = template_slide.shapes[rect_shape_idx]
#         except Exception:
#             rect_shape = None

#     if rect_shape is None:
#         rect_shape = _find_rectangle_shape(template_slide)

#     if rect_shape is None:
#         shapes_info = []
#         for i, sh in enumerate(template_slide.shapes):
#             shapes_info.append(f"#{i} name={getattr(sh,'name',None)} type={getattr(sh,'shape_type',None)} w={getattr(sh,'width',None)} h={getattr(sh,'height',None)}")
#         raise ValueError("Could not find a rectangle shape to replace on template slide. Slide shapes:\n" + "\n".join(shapes_info))

#     # capture geometry and remove the rectangle
#     left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height
#     _safe_remove_shape(template_slide, rect_shape)

#     # create a single table that fits within the rectangle; reduce row height so more rows fit visually
#     def _add_table_on_slide(slide, df, rows_total):
#         cols = max(1, len(df.columns))
#         tbl_sh = slide.shapes.add_table(rows=rows_total, cols=cols,
#                                        left=left, top=top, width=width, height=height)
#         tbl = tbl_sh.table
#         # header row
#         for c, col in enumerate(df.columns):
#             if c >= len(tbl.columns):
#                 break
#             tbl.rows[0].cells[c].text = str(col)
#         # equal column widths
#         col_width = width // max(1, cols)
#         for c in range(min(cols, len(tbl.columns))):
#             tbl.columns[c].width = col_width
#         # reduce row heights: compute a smaller per-row height but keep a minimum
#         try:
#             # rows_total includes header; compute target row height
#             min_row_height = int(Pt(10)) if 'Pt' in globals() else int(91440)  # fallback if Pt not imported
#             target_row_height = max(int(height // max(2, rows_total)), min_row_height)
#             for r in range(min(rows_total, len(tbl.rows))):
#                 try:
#                     tbl.rows[r].height = target_row_height
#                 except Exception:
#                     # some pptx versions ignore row.height; ignore failures
#                     pass
#         except Exception:
#             pass
#         return tbl_sh

#     remaining_df = sheet_df.reset_index(drop=True).copy()

#     # Determine how many rows we can write on this single slide (header + data rows)
#     rows_to_create = max(2, min(max_rows_per_slide, len(remaining_df) + 1))
#     table_shape = _add_table_on_slide(template_slide, remaining_df, rows_to_create)

#     tbl = table_shape.table
#     capacity = len(tbl.rows) - 1  # data rows capacity (row 0 is header)

#     # write only up to capacity rows and do not create duplicates
#     chunk = remaining_df.iloc[:capacity].reset_index(drop=True)
#     rows_written = 0

#     for r_idx, (_, row) in enumerate(chunk.iterrows()):
#         if r_idx >= capacity:
#             break
#         for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
#             try:
#                 cell = tbl.rows[1 + r_idx].cells[c_idx]
#             except Exception:
#                 continue
#             overflow, first_part, remaining = _split_cell_text_by_lines(val, max_lines_per_cell)
#             # write only the first part that fits; do not attempt to carry remaining text to new slides
#             cell.text = first_part or ""
#         rows_written += 1

#     # done — we intentionally do not process remaining rows or create duplicates
#     print(f"Wrote {rows_written} row(s) to slide {template_index + 1}. Remaining rows were not added (no duplication).")

def replace_rectangle_with_table_and_handle_overflow(
        ppt,
        template_slide,               # expects a pptx.slide.Slide object
        sheet_df: pd.DataFrame,
        rect_shape_idx=None,
        max_rows_per_slide=20,
        max_lines_per_cell=14):
    """
    Simple, clear implementation that:
      - Inserts a table into the largest rectangle on template_slide.
      - Uses the first row of sheet_df as the table header and remaining rows as data.
      - Reduces row height so the table fills the rectangle.
      - If data rows exceed capacity, duplicates the current slide (using existing duplicator)
        and continues writing remaining rows on the next slide.
      - Uses Arial 10.5 for all table text.
    Assumes helper functions already present in the module:
      _find_rectangle_shape, _safe_remove_shape, _split_cell_text_by_lines,
      duplicate_slide_safe_insert_after
    """
    if sheet_df is None or len(sheet_df) == 0:
        return 0

    # normalize dataframe
    df = sheet_df.reset_index(drop=True).copy()
    # treat first row as header, rest as data
    header = ["" if v is None else str(v) for v in df.iloc[0].tolist()]
    data = df.iloc[1:].reset_index(drop=True)

    # find rectangle on template slide
    rect = None
    if rect_shape_idx is not None:
        rect = template_slide.shapes[rect_shape_idx]
    if rect is None:
        rect = _find_rectangle_shape(template_slide)
    if rect is None:
        raise ValueError("Rectangle shape not found on template slide.")

    left, top, width, height = rect.left, rect.top, rect.width, rect.height
    _safe_remove_shape(template_slide, rect)

    # capacity: data rows per slide (header occupies one row)
    data_rows_per_slide = max(1, max_rows_per_slide - 1)
    total_written = 0
    current_slide = template_slide

    # iterate over data in chunks that fit per slide
    start = 0
    while start < len(data):
        chunk = data.iloc[start:start + data_rows_per_slide].reset_index(drop=True)
        rows_needed = 1 + len(chunk)            # header + data rows
        cols = max(1, len(header))

        # add table to current slide sized to rectangle
        tbl_sh = current_slide.shapes.add_table(rows=rows_needed, cols=cols,
                                                left=left, top=top, width=width, height=height)
        tbl = tbl_sh.table

        # set equal column widths
        col_width = width // max(1, cols)
        for c in range(cols):
            tbl.columns[c].width = col_width

        # set row heights so table fills rectangle (simple integer division)
        row_height = int(height // max(1, rows_needed))
        for r in range(rows_needed):
            try:
                tbl.rows[r].height = row_height
            except Exception:
                pass

        # write header (Arial 10.5 bold)
        for c, hval in enumerate(header[:cols]):
            cell = tbl.rows[0].cells[c]
            cell.text = "" if hval is None else str(hval)
            p = cell.text_frame.paragraphs[0]
            if not p.runs:
                run = p.add_run()
            else:
                run = p.runs[0]
            run.text = cell.text_frame.text
            run.font.name = "Arial"
            run.font.size = Pt(10.5)
            run.font.bold = True

        # write data rows (Arial 10.5). enforce max_lines_per_cell per cell.
        for r_idx, (_, row) in enumerate(chunk.iterrows()):
            for c_idx, val in enumerate(row.tolist()[:cols]):
                cell = tbl.rows[1 + r_idx].cells[c_idx]
                text = "" if val is None else str(val)
                lines = text.splitlines()
                if len(lines) > max_lines_per_cell:
                    text_to_write = "\n".join(lines[:max_lines_per_cell])
                else:
                    text_to_write = text
                cell.text = text_to_write
                p = cell.text_frame.paragraphs[0]
                if not p.runs:
                    run = p.add_run()
                else:
                    run = p.runs[0]
                run.text = text_to_write
                run.font.name = "Arial"
                run.font.size = Pt(10.5)
                run.font.bold = False

        total_written += len(chunk)
        start += len(chunk)

        # if more rows remain, duplicate current_slide and set duplicate as next current_slide
        if start < len(data):
            current_slide = duplicate_slide_safe_insert_after(ppt, current_slide)

    return total_written


from pptx.util import Pt
# #for the Annexture Table
# def ann_replace_rectangle_with_table_and_handle_overflow(
#         ppt,
#         template_slide,               # expects a pptx.slide.Slide object
#         sheet_df: pd.DataFrame,
#         rect_shape_idx: Optional[int] = None,
#         max_rows_per_slide: int = 20,
#         max_lines_per_cell: int = 14):
#     """
#     Insert a single table into the largest rectangle on template_slide using:
#       - A1 -> replaces any shape text exactly "{{Annexure-Title}}" (Arial, 10.5pt, bold)
#       - header row taken from sheet row 3 (index 2)
#       - data rows start from sheet row 4 (index 3)
#       - NO slide duplication: write only the maximum rows that fit on one slide
#       - table height is exactly the rectangle height; row heights are reduced to fit
#     """
#     if sheet_df is None or len(sheet_df) == 0:
#         print("No data in sheet_df; nothing to write.")
#         return
#     template_slide1 = template_slide
#     # --- read title from A1 (safe) ---
#     df_reset_for_title = sheet_df.reset_index(drop=True)
#     title_value = ""
#     try:
#         if len(df_reset_for_title) > 0:
#             # try iat for speed; fallback to iloc
#             try:
#                 title_value = df_reset_for_title.iat[0, 0]
#             except Exception:
#                 title_value = df_reset_for_title.iloc[0, 0]
#             if title_value is None:
#                 title_value = ""
#             title_value = str(title_value)
#     except Exception:
#         title_value = ""

#     # replace token "{{Annexure-Title}}" on template_slide with title_value (Arial 10.5pt, bold)
#     if title_value:
#         for sh in template_slide1.shapes:
#             try:
#                 tf = getattr(sh, "text_frame", None)
#                 if tf is None:
#                     continue
#                 cur = (tf.text or "").strip()
#                 if cur == "{{Annexure-Title}}":
#                     # clear and set new text with formatting
#                     tf.clear()
#                     p = tf.paragraphs[0]
#                     run = p.add_run()
#                     run.text = title_value
#                     run.font.name = "Arial"
#                     run.font.size = Pt(10.5)
#                     run.font.bold = True
#                     break
#             except Exception:
#                 # ignore shape-level errors and continue
#                 continue

#     # --- determine header (row 3) and data (rows from 4 onwards) ---
#     df_reset = sheet_df.reset_index(drop=True).copy()
#     if len(df_reset) > 2:
#         header_row = df_reset.iloc[2]
#         header_values = [header_row.iloc[c] if c < len(header_row) else "" for c in range(len(df_reset.columns))]
#     else:
#         header_values = list(df_reset.columns)

#     data_df = df_reset.iloc[3:].reset_index(drop=True).copy()  # rows from index 3 onward

#     # --- find rectangle shape (use existing helper) ---
#     rect_shape = None
#     if rect_shape_idx is not None:
#         try:
#             rect_shape = template_slide1.shapes[rect_shape_idx]
#         except Exception:
#             rect_shape = None

#     if rect_shape is None:
#         rect_shape = _find_rectangle_shape(template_slide1)

#     if rect_shape is None:
#         shapes_info = []
#         for i, sh in enumerate(template_slide1.shapes):
#             shapes_info.append(f"#{i} name={getattr(sh,'name',None)} type={getattr(sh,'shape_type',None)} w={getattr(sh,'width',None)} h={getattr(sh,'height',None)}")
#         raise ValueError("Could not find a rectangle shape to replace on template slide. Slide shapes:\n" + "\n".join(shapes_info))

#     # capture geometry and remove the rectangle
#     left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height
#     _safe_remove_shape(template_slide1, rect_shape)

#     # --- compute rows to create (header + data rows that fit) ---
#     # We will create at most max_rows_per_slide rows (including header).
#     # rows_available_data = how many data rows we can show = rows_to_create - 1
#     # choose rows_to_create = min(max_rows_per_slide, 1 + len(data_df))
#     rows_to_create = max(2, min(max_rows_per_slide, 1 + len(data_df)))
#     cols_count = max(1, len(df_reset.columns))

#     # --- add table exactly to rectangle bounds ---
#     tbl_sh = template_slide1.shapes.add_table(rows=rows_to_create, cols=cols_count,
#                                              left=left, top=top, width=width, height=height)
#     tbl = tbl_sh.table

#     # set header row from header_values (use only up to cols_count)
#     for c, col_val in enumerate(header_values[:cols_count]):
#         try:
#             cell = tbl.rows[0].cells[c]
#             cell.text = "" if col_val is None else str(col_val)
#             # format header cell text run (optional: keep simple)
#             try:
#                 p = cell.text_frame.paragraphs[0]
#                 if not p.runs:
#                     r = p.add_run()
#                     r.text = cell.text_frame.text
#                 else:
#                     r = p.runs[0]
#                 r.font.name = "Arial"
#                 r.font.size = Pt(10.5)
#                 r.font.bold = True
#             except Exception:
#                 pass
#         except Exception:
#             continue

#     # set column widths equally (best-effort)
#     try:
#         col_width = width // max(1, cols_count)
#         for c in range(min(cols_count, len(tbl.columns))):
#             tbl.columns[c].width = col_width
#     except Exception:
#         pass

#     # set row heights so table exactly fills rectangle height and rows are reduced to fit
#     try:
#         rows_total_safe = max(1, rows_to_create)
#         target_row_height = max(int(height // rows_total_safe), int(Pt(6)))  # minimum ~6pt
#         for r in range(min(rows_to_create, len(tbl.rows))):
#             try:
#                 tbl.rows[r].height = target_row_height
#             except Exception:
#                 pass
#     except Exception:
#         pass

#     # --- write data rows up to capacity (no duplication, no overflow handling beyond first part) ---
#     capacity = len(tbl.rows) - 1  # number of data rows available
#     chunk = data_df.iloc[:capacity].reset_index(drop=True)
#     rows_written = 0

#     for r_idx, (_, row) in enumerate(chunk.iterrows()):
#         if r_idx >= capacity:
#             break
#         for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
#             try:
#                 cell = tbl.rows[1 + r_idx].cells[c_idx]
#             except Exception:
#                 continue
#             # split by lines but only keep first part (no overflow slides)
#             overflow, first_part, remaining = _split_cell_text_by_lines(val, max_lines_per_cell)
#             text_to_write = first_part or ""
#             # assign text and apply a readable font (Arial, 10pt)
#             try:
#                 cell.text = text_to_write
#                 p = cell.text_frame.paragraphs[0]
#                 if not p.runs:
#                     r = p.add_run()
#                     r.text = text_to_write
#                 else:
#                     r = p.runs[0]
#                 r.font.name = "Arial"
#                 r.font.size = Pt(10)
#                 r.font.bold = False
#             except Exception:
#                 # fallback: direct assignment
#                 try:
#                     cell.text = text_to_write
#                 except Exception:
#                     pass
#         rows_written += 1

#     # done: intentionally do not process remaining rows or create duplicates
#     template_index = list(ppt.slides).index(template_slide1)
#     print(f"Wrote {rows_written} data row(s) to slide {template_index + 1}. Remaining rows were not added (no duplication).")


from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from typing import Optional
import pandas as pd
from copy import deepcopy
import io

# def replace_rectangle_with_table_and_handle_overflow(
#         ppt,
#         template_slide,               # expects a pptx.slide.Slide object
#         sheet_df: pd.DataFrame,
#         rect_shape_idx: Optional[int] = None,
#         max_rows_per_slide: int = 20,
#         max_lines_per_cell: int = 14):
#     """
#     Replace the largest rectangle on template_slide with a table built from sheet_df.
#     If the data overflows the capacity of one slide, duplicate the template slide and
#     continue writing the remaining rows across additional slides until all rows are written.
#     Duplication is implemented by recreating shapes (textboxes, tables, pictures) to avoid
#     duplicate relationship IDs when saving.
#     """

#     if sheet_df is None or len(sheet_df) == 0:
#         print("No data in sheet_df; nothing to write.")
#         return

#     # ----------------- helpers -----------------
#     def _find_rectangle_shape(slide):
#         """Find a rectangle-like shape to replace (best-effort)."""
#         best = None
#         best_area = 0
#         for sh in slide.shapes:
#             try:
#                 st = getattr(sh, "shape_type", None)
#                 w = getattr(sh, "width", 0)
#                 h = getattr(sh, "height", 0)
#                 area = int(w) * int(h)
#                 # treat text boxes and auto shapes as candidates
#                 if st in (1, 17, 19) or getattr(sh, "has_text_frame", False):
#                     if area > best_area:
#                         best = sh
#                         best_area = area
#             except Exception:
#                 continue
#         return best

#     def _safe_remove_shape(slide, shape):
#         """Remove shape from slide safely."""
#         try:
#             slide.shapes._spTree.remove(shape.element)
#         except Exception:
#             try:
#                 slide.shapes._spTree.remove(shape._element)
#             except Exception:
#                 pass

#     def _split_cell_text_by_lines(val, max_lines):
#         """Return (overflow_flag, first_part, remaining). Works on strings or numbers."""
#         if val is None:
#             return False, "", ""
#         s = str(val)
#         lines = s.splitlines() or [s]
#         if len(lines) <= max_lines:
#             return False, s, ""
#         first = "\n".join(lines[:max_lines])
#         rem = "\n".join(lines[max_lines:])
#         return True, first, rem

#     def _add_table_on_slide(slide, left, top, width, height, header_values, cols_count, rows_total):
#         """Add a table exactly to the rectangle bounds and set header text."""
#         tbl_sh = slide.shapes.add_table(rows=rows_total, cols=cols_count,
#                                        left=left, top=top, width=width, height=height)
#         tbl = tbl_sh.table
#         # header row: set header_values (only non-empty)
#         for c in range(min(cols_count, len(header_values))):
#             hv = header_values[c] if c < len(header_values) else ""
#             if hv is None:
#                 hv = ""
#             try:
#                 cell = tbl.rows[0].cells[c]
#                 if str(hv).strip() != "":
#                     try:
#                         cell.text_frame.clear()
#                     except Exception:
#                         pass
#                     cell.text = str(hv)
#                     p = cell.text_frame.paragraphs[0]
#                     if not p.runs:
#                         rrun = p.add_run()
#                         rrun.text = str(hv)
#                     else:
#                         rrun = p.runs[0]
#                         rrun.text = str(hv)
#                     rrun.font.name = "Arial"
#                     rrun.font.size = Pt(10.5)
#                     rrun.font.bold = True
#             except Exception:
#                 continue
#         # equal column widths (best-effort)
#         try:
#             col_width = width // max(1, cols_count)
#             for c in range(min(cols_count, len(tbl.columns))):
#                 tbl.columns[c].width = col_width
#         except Exception:
#             pass
#         # set row heights so table exactly fills rectangle height and rows are reduced to fit
#         try:
#             rows_total_safe = max(1, rows_total)
#             target_row_height = max(int(height // rows_total_safe), int(Pt(6)))
#             for r in range(min(rows_total, len(tbl.rows))):
#                 try:
#                     tbl.rows[r].height = target_row_height
#                 except Exception:
#                     pass
#         except Exception:
#             pass
#         return tbl_sh

#     # ----------------- prepare data -----------------

#     df_reset = sheet_df.reset_index(drop=True).copy()

#     # header row from sheet row 3 (index 2) if present, else use df columns
#     if len(df_reset) > 2:
#         header_row = df_reset.iloc[2]
#         header_values = [header_row.iloc[c] if c < len(header_row) else "" for c in range(len(df_reset.columns))]
#     else:
#         header_values = list(df_reset.columns)

#     # data rows start from sheet row 4 (index 3)
#     data_df = df_reset.iloc[3:].reset_index(drop=True).copy()

#     # ----------------- locate rectangle -----------------

#     rect_shape = None
#     if rect_shape_idx is not None:
#         try:
#             rect_shape = template_slide.shapes[rect_shape_idx]
#         except Exception:
#             rect_shape = None

#     if rect_shape is None:
#         rect_shape = _find_rectangle_shape(template_slide)

#     if rect_shape is None:
#         shapes_info = []
#         for i, sh in enumerate(template_slide.shapes):
#             shapes_info.append(f"#{i} name={getattr(sh,'name',None)} type={getattr(sh,'shape_type',None)} w={getattr(sh,'width',None)} h={getattr(sh,'height',None)}")
#         raise ValueError("Could not find a rectangle shape to replace on template slide. Slide shapes:\n" + "\n".join(shapes_info))

#     left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height

#     # ----------------- compute capacities -----------------

#     data_rows_per_slide = max(1, max_rows_per_slide - 1)
#     total_data_rows = len(data_df)
#     if total_data_rows == 0:
#         _safe_remove_shape(template_slide, rect_shape)
#         _add_table_on_slide(template_slide, left, top, width, height, header_values, max(1, len(df_reset.columns)), 1)
#         print("No data rows to write; created header-only table.")
#         return

#     # ----------------- write across slides -----------------

#     remaining = data_df.copy()
#     slide_to_use = template_slide
#     first_iteration = True
#     rows_written_total = 0

#     while len(remaining) > 0:
#         if first_iteration:
#             # remove rectangle from original template slide
#             _safe_remove_shape(slide_to_use, rect_shape)
#             first_iteration = False
#         else:
#             # create a safe duplicate by recreating shapes
#             slide_to_use = duplicate_slide(ppt, template_slide)
#             # remove rectangle on duplicated slide if found
#             dup_rect = _find_rectangle_shape(slide_to_use)
#             if dup_rect is not None:
#                 _safe_remove_shape(slide_to_use, dup_rect)

#         rows_for_this_slide = min(max_rows_per_slide, 1 + len(remaining))
#         rows_for_this_slide = max(2, rows_for_this_slide)
#         cols_count = max(1, len(df_reset.columns))

#         table_shape = _add_table_on_slide(slide_to_use, left, top, width, height, header_values, cols_count, rows_for_this_slide)
#         tbl = table_shape.table
#         capacity = len(tbl.rows) - 1

#         chunk = remaining.iloc[:capacity].reset_index(drop=True)
#         for r_idx, (_, row) in enumerate(chunk.iterrows()):
#             for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
#                 try:
#                     cell = tbl.rows[1 + r_idx].cells[c_idx]
#                 except Exception:
#                     continue
#                 overflow, first_part, _rem = _split_cell_text_by_lines(val, max_lines_per_cell)
#                 text_to_write = first_part or ""
#                 try:
#                     try:
#                         cell.text_frame.clear()
#                     except Exception:
#                         pass
#                     cell.text = text_to_write
#                     p = cell.text_frame.paragraphs[0]
#                     if not p.runs:
#                         rrun = p.add_run()
#                         rrun.text = text_to_write
#                     else:
#                         rrun = p.runs[0]
#                         rrun.text = text_to_write
#                     rrun.font.name = "Arial"
#                     rrun.font.size = Pt(10)
#                     rrun.font.bold = False
#                 except Exception:
#                     try:
#                         cell.text = text_to_write
#                     except Exception:
#                         pass

#         rows_written = len(chunk)
#         rows_written_total += rows_written
#         remaining = remaining.iloc[rows_written:].reset_index(drop=True)

#     template_index = list(ppt.slides).index(template_slide)
#     print(f"Wrote {rows_written_total} data row(s) across slides starting at template slide index {template_index + 1}.")

from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from typing import Optional
import pandas as pd

from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def duplicate_slide_safe_insert_after(pres, src_slide):
    """
    Create a new slide using the same layout as src_slide and recreate shapes
    (text frames, tables, pictures) on the new slide using python-pptx API.
    Place the newly created slide immediately after src_slide in the slide order.
    Returns the new slide object.
    """
    # create new slide (this appends at the end)
    new_slide = pres.slides.add_slide(src_slide.slide_layout)

    # recreate shapes (text, tables, pictures) on the new slide
    for sh in src_slide.shapes:
        try:
            stype = sh.shape_type
        except Exception:
            stype = None

        # Text-containing shapes (text boxes, placeholders, auto shapes)
        if getattr(sh, "has_text_frame", False):
            left, top, width, height = sh.left, sh.top, sh.width, sh.height
            tx = new_slide.shapes.add_textbox(left, top, width, height)
            src_tf = sh.text_frame
            dst_tf = tx.text_frame
            dst_tf.clear()
            for p in src_tf.paragraphs:
                new_p = dst_tf.add_paragraph()
                new_p.level = p.level
                for run in p.runs:
                    new_run = new_p.add_run()
                    new_run.text = run.text
                    try:
                        new_run.font.name = run.font.name
                        new_run.font.size = run.font.size
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic
                        if run.font.color and run.font.color.rgb:
                            new_run.font.color.rgb = run.font.color.rgb
                    except Exception:
                        pass
            # remove initial empty paragraph if present
            try:
                if dst_tf.paragraphs and dst_tf.paragraphs[0].text == "":
                    dst_tf.paragraphs.pop(0)
            except Exception:
                pass

        # Tables
        elif stype == MSO_SHAPE_TYPE.TABLE:
            left, top, width, height = sh.left, sh.top, sh.width, sh.height
            tbl = sh.table
            rows = len(tbl.rows)
            cols = len(tbl.columns)
            new_tbl_sh = new_slide.shapes.add_table(rows=rows, cols=cols,
                                                    left=left, top=top, width=width, height=height)
            new_tbl = new_tbl_sh.table
            for r in range(rows):
                for c in range(cols):
                    try:
                        src_cell = tbl.rows[r].cells[c]
                        dst_cell = new_tbl.rows[r].cells[c]
                        text = src_cell.text_frame.text if getattr(src_cell, "text_frame", None) else ""
                        try:
                            dst_cell.text_frame.clear()
                        except Exception:
                            pass
                        dst_cell.text = text
                        # copy first run formatting if present
                        try:
                            src_p = src_cell.text_frame.paragraphs[0]
                            if src_p.runs:
                                src_run = src_p.runs[0]
                                dst_p = dst_cell.text_frame.paragraphs[0]
                                if not dst_p.runs:
                                    rrun = dst_p.add_run()
                                else:
                                    rrun = dst_p.runs[0]
                                rrun.text = src_run.text
                                try:
                                    rrun.font.name = src_run.font.name
                                    rrun.font.size = src_run.font.size
                                    rrun.font.bold = src_run.font.bold
                                    if src_run.font.color and src_run.font.color.rgb:
                                        rrun.font.color.rgb = src_run.font.color.rgb
                                except Exception:
                                    pass
                        except Exception:
                            pass
                        # copy fill if present (best-effort)
                        try:
                            if src_cell.fill.type is not None:
                                dst_cell.fill.solid()
                                dst_cell.fill.fore_color.rgb = src_cell.fill.fore_color.rgb
                        except Exception:
                            pass
                    except Exception:
                        continue

        # Pictures
        elif stype == MSO_SHAPE_TYPE.PICTURE:
            try:
                img = sh.image
                image_bytes = img.blob
                left, top, width, height = sh.left, sh.top, sh.width, sh.height
                new_slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width, height=height)
            except Exception:
                continue

        # Other shapes with text: recreate as textbox
        else:
            if getattr(sh, "has_text_frame", False):
                left, top, width, height = sh.left, sh.top, sh.width, sh.height
                tx = new_slide.shapes.add_textbox(left, top, width, height)
                try:
                    tx.text = sh.text
                except Exception:
                    pass
            else:
                # skip complex shapes to avoid relationship duplication
                continue

    # Move the newly created slide to be immediately after src_slide
    try:
        sldIdLst = pres.slides._sldIdLst  # internal slide id list
        # the new slide's sldId is the last element
        new_sldId = sldIdLst[-1]
        # find index of src_slide in presentation
        src_index = list(pres.slides).index(src_slide)
        # remove the last element and insert it after src_index
        sldIdLst.remove(new_sldId)
        sldIdLst.insert(src_index + 1, new_sldId)
    except Exception:
        # if anything goes wrong with reordering, leave the slide appended at the end
        pass

    return new_slide

def ann_replace_rectangle_with_table_and_handle_overflow(
        ppt,
        template_slide,               # expects a pptx.slide.Slide object
        sheet_df: pd.DataFrame,
        rect_shape_idx: Optional[int] = None,
        max_rows_per_slide: int = 20,
        max_lines_per_cell: int = 14):
    """
    Same logic as before but operate on a duplicate of template_slide so the original
    template can be reused multiple times without being overwritten.
    """
    if sheet_df is None or len(sheet_df) == 0:
        print("No data in sheet_df; nothing to write.")
        return

    # duplicate template slide safely and operate on the duplicate
    working_slide = duplicate_slide_safe_insert_after(ppt, template_slide)

    # rest of your original function logic, replacing template_slide1 with working_slide
    template_slide1 = working_slide

    # --- read title from A1 (safe) ---
    df_reset_for_title = sheet_df.reset_index(drop=True)
    title_value = ""
    try:
        if len(df_reset_for_title) > 0:
            try:
                title_value = df_reset_for_title.iat[0, 0]
            except Exception:
                title_value = df_reset_for_title.iloc[0, 0]
            if title_value is None:
                title_value = ""
            title_value = str(title_value)
    except Exception:
        title_value = ""

    if title_value:
        for sh in template_slide1.shapes:
            try:
                tf = getattr(sh, "text_frame", None)
                if tf is None:
                    continue
                cur = (tf.text or "").strip()
                if cur == "{{Annexure-Title}}":
                    tf.clear()
                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = title_value
                    run.font.name = "Arial"
                    run.font.size = Pt(24)
                    run.font.bold = True
                    break
            except Exception:
                continue

    # --- determine header and data ---
    df_reset = sheet_df.reset_index(drop=True).copy()
    if len(df_reset) > 2:
        header_row = df_reset.iloc[2]
        header_values = [header_row.iloc[c] if c < len(header_row) else "" for c in range(len(df_reset.columns))]
    else:
        header_values = list(df_reset.columns)

    data_df = df_reset.iloc[3:].reset_index(drop=True).copy()

    # --- find rectangle shape on the duplicated slide ---
    rect_shape = None
    if rect_shape_idx is not None:
        try:
            rect_shape = template_slide1.shapes[rect_shape_idx]
        except Exception:
            rect_shape = None

    if rect_shape is None:
        rect_shape = _find_rectangle_shape(template_slide1)

    if rect_shape is None:
        shapes_info = []
        for i, sh in enumerate(template_slide1.shapes):
            shapes_info.append(f"#{i} name={getattr(sh,'name',None)} type={getattr(sh,'shape_type',None)} w={getattr(sh,'width',None)} h={getattr(sh,'height',None)}")
        raise ValueError("Could not find a rectangle shape to replace on template slide. Slide shapes:\n" + "\n".join(shapes_info))

    left, top, width, height = rect_shape.left, rect_shape.top, rect_shape.width, rect_shape.height
    _safe_remove_shape(template_slide1, rect_shape)

    rows_to_create = max(2, min(max_rows_per_slide, 1 + len(data_df)))
    cols_count = max(1, len(df_reset.columns))

    tbl_sh = template_slide1.shapes.add_table(rows=rows_to_create, cols=cols_count,
                                             left=left, top=top, width=width, height=height)
    tbl = tbl_sh.table

    for c, col_val in enumerate(header_values[:cols_count]):
        try:
            cell = tbl.rows[0].cells[c]
            cell.text = "" if col_val is None else str(col_val)
            try:
                p = cell.text_frame.paragraphs[0]
                if not p.runs:
                    r = p.add_run()
                    r.text = cell.text_frame.text
                else:
                    r = p.runs[0]
                r.font.name = "Arial"
                r.font.size = Pt(10.5)
                r.font.bold = True
            except Exception:
                pass
        except Exception:
            continue

    try:
        col_width = width // max(1, cols_count)
        for c in range(min(cols_count, len(tbl.columns))):
            tbl.columns[c].width = col_width
    except Exception:
        pass

    try:
        rows_total_safe = max(1, rows_to_create)
        target_row_height = max(int(height // rows_total_safe), int(Pt(6)))
        for r in range(min(rows_to_create, len(tbl.rows))):
            try:
                tbl.rows[r].height = target_row_height
            except Exception:
                pass
    except Exception:
        pass

    capacity = len(tbl.rows) - 1
    chunk = data_df.iloc[:capacity].reset_index(drop=True)
    rows_written = 0

    for r_idx, (_, row) in enumerate(chunk.iterrows()):
        if r_idx >= capacity:
            break
        for c_idx, val in enumerate(row.tolist()[:len(tbl.columns)]):
            try:
                cell = tbl.rows[1 + r_idx].cells[c_idx]
            except Exception:
                continue
            overflow, first_part, remaining = _split_cell_text_by_lines(val, max_lines_per_cell)
            text_to_write = first_part or ""
            try:
                cell.text = text_to_write
                p = cell.text_frame.paragraphs[0]
                if not p.runs:
                    r = p.add_run()
                    r.text = text_to_write
                else:
                    r = p.runs[0]
                r.font.name = "Arial"
                r.font.size = Pt(10)
                r.font.bold = False
            except Exception:
                try:
                    cell.text = text_to_write
                except Exception:
                    pass
        rows_written += 1

    template_index = list(ppt.slides).index(template_slide1)
    print(f"Wrote {rows_written} data row(s) to slide {template_index + 1}. (This is a duplicate of the template slide; original template remains unchanged.)")

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import pandas as pd

def apply_table_style(tbl):
    """Apply required visual style to a pptx table object (header fill, fonts, alignment, spacing, borders)."""
    header_fill = RGBColor(0, 51, 141)
    border_color = RGBColor(128, 128, 128)
    font_name = "Arial"
    font_size = Pt(10.5)
    # helper to set cell border (all four sides)
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls

    def _set_cell_border(cell, rgb=(128,128,128), width_emu=12700):
        r, g, b = rgb
        val = "{:02X}{:02X}{:02X}".format(r, g, b)
        xml = f"""
        <a:tcPr {nsdecls('a')}>
        <a:lnL w="{width_emu}"><a:solidFill><a:srgbClr val="{val}"/></a:solidFill></a:lnL>
        <a:lnR w="{width_emu}"><a:solidFill><a:srgbClr val="{val}"/></a:solidFill></a:lnR>
        <a:lnT w="{width_emu}"><a:solidFill><a:srgbClr val="{val}"/></a:solidFill></a:lnT>
        <a:lnB w="{width_emu}"><a:solidFill><a:srgbClr val="{val}"/></a:solidFill></a:lnB>
        </a:tcPr>
        """
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        new_tcPr = parse_xml(xml)
        # append each child of new_tcPr into tcPr
        for child in new_tcPr:
            tcPr.append(child)

    rows = len(tbl.rows)
    cols = len(tbl.columns)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.rows[r].cells[c]
            text = (cell.text or "")
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            p.line_spacing = 1.2
            if not p.runs:
                run = p.add_run()
            else:
                run = p.runs[0]
            run.text = p.text
            run.font.name = font_name
            run.font.size = font_size
            run.font.bold = (r == 0)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            _set_cell_border(cell)

def process_key_stats_sheets(ppt, df_sheets, template_index=9, rect_shape_idx=None, max_rows_per_slide=20):
    """
    For each Excel sheet whose name starts with 'Key stats' create one duplicate of template slide (index template_index),
    replace the rectangle there with a table built from the sheet, and style it with apply_table_style.
    Returns list of slide_ids created.
    """
    prs = ppt
    sheets = df_sheets
    key_sheets = {name: df for name, df in sheets.items() if name.lower().startswith("key stats")}
    created_slide_ids = []

    def duplicate_slide_after(pres, src_slide):
        new = pres.slides.add_slide(src_slide.slide_layout)
        for sh in src_slide.shapes:
            if getattr(sh, "has_text_frame", False):
                left, top, w, h = sh.left, sh.top, sh.width, sh.height
                tx = new.shapes.add_textbox(left, top, w, h)
                dst_tf = tx.text_frame
                dst_tf.clear()
                for p in sh.text_frame.paragraphs:
                    np = dst_tf.add_paragraph()
                    np.level = p.level
                    for run in p.runs:
                        nr = np.add_run()
                        nr.text = run.text
                        if run.font is not None:
                            if run.font.name: nr.font.name = run.font.name
                            if run.font.size: nr.font.size = run.font.size
                            nr.font.bold = bool(run.font.bold)
                            if getattr(run.font, "color", None) and getattr(run.font.color, "rgb", None):
                                nr.font.color.rgb = run.font.color.rgb
            elif getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
                left, top, w, h = sh.left, sh.top, sh.width, sh.height
                tbl = sh.table
                rows, cols = len(tbl.rows), len(tbl.columns)
                new_tbl_sh = new.shapes.add_table(rows=rows, cols=cols, left=left, top=top, width=w, height=h)
                new_tbl = new_tbl_sh.table
                for r in range(rows):
                    for c in range(cols):
                        src_cell = tbl.rows[r].cells[c]
                        dst_cell = new_tbl.rows[r].cells[c]
                        dst_cell.text = src_cell.text_frame.text if getattr(src_cell, "text_frame", None) else ""
            elif getattr(sh, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = sh.image
                    new.shapes.add_picture(img.blob, sh.left, sh.top, width=sh.width, height=sh.height)
                except:
                    pass
        # move new slide immediately after src_slide
        sldIdLst = pres.slides._sldIdLst
        new_id = sldIdLst[-1]
        src_idx = list(pres.slides).index(src_slide)
        sldIdLst.remove(new_id)
        sldIdLst.insert(src_idx + 1, new_id)
        return pres.slides[src_idx + 1]

    def find_rectangle(slide):
        best = None
        best_area = 0
        for sh in slide.shapes:
            w = getattr(sh, "width", 0)
            h = getattr(sh, "height", 0)
            area = int(w) * int(h)
            if getattr(sh, "has_text_frame", False) or getattr(sh, "shape_type", None) in (1, 17, 19):
                if area > best_area:
                    best = sh
                    best_area = area
        return best

    template_slide = prs.slides[template_index]
    for name, df in key_sheets.items():
        slide = duplicate_slide_after(prs, template_slide)
        # find and remove rectangle
        rect = None
        if rect_shape_idx is not None:
            rect = slide.shapes[rect_shape_idx]
        if rect is None:
            rect = find_rectangle(slide)
        left, top, width, height = rect.left, rect.top, rect.width, rect.height
        slide.shapes._spTree.remove(rect._element)
        # create table (header = first row if header-like else use df.columns)
        if df.shape[0] > 0 and all(str(x).strip() != "" for x in df.iloc[0].tolist()):
            header = [str(x) for x in df.iloc[0].tolist()]
            data = df.iloc[1:].reset_index(drop=True)
        else:
            header = [str(c) for c in df.columns]
            data = df.reset_index(drop=True)
        rows_needed = 1 + len(data)
        cols = max(1, len(header))
        tbl_sh = slide.shapes.add_table(rows=rows_needed, cols=cols, left=left, top=top, width=width, height=height)
        tbl = tbl_sh.table
        col_w = width // max(1, cols)
        for c in range(cols):
            tbl.columns[c].width = col_w
        row_h = int(height // max(1, rows_needed))
        for r in range(rows_needed):
            tbl.rows[r].height = row_h
        # fill header
        for c, h in enumerate(header[:cols]):
            cell = tbl.rows[0].cells[c]
            cell.text = "" if h is None else str(h)
        # fill data
        for r_idx in range(len(data)):
            row = data.iloc[r_idx].tolist()
            for c_idx in range(min(len(row), cols)):
                cell = tbl.rows[1 + r_idx].cells[c_idx]
                cell.text = "" if row[c_idx] is None else str(row[c_idx])
        apply_table_style(tbl)
        created_slide_ids.append(slide.slide_id)
    return created_slide_ids

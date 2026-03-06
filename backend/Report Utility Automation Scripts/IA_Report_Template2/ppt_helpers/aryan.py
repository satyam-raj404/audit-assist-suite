import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml import parse_xml
import copy
from copy import deepcopy
from pptx.oxml.ns import qn
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def _iter_table_cells_in_shape(shape):
    """
    Yield every table cell reachable from this shape, including:
    - Direct table shapes
    - Tables nested inside group shapes (recursive)
    """
    # If it's a group, recurse into children
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shp in shape.shapes:
            yield from _iter_table_cells_in_shape(shp)
        return

    # If the shape has a table, iterate its cells
    if getattr(shape, "has_table", False) and shape.has_table:
        tbl = shape.table
        # Note: merged cells are repeated in the grid; operating on all is safe for text/formatting.
        for r in tbl.rows:
            for c in r.cells:
                yield c

def _text_matches(haystack: str, needle: str, *, case_sensitive: bool, exact_match: bool) -> bool:
    if not case_sensitive:
        haystack, needle = haystack.lower(), needle.lower()
    return (haystack == needle) if exact_match else (needle in haystack)

def _replace_cell_text_and_format(cell, replacement_text: str):
    """
    Replace entire cell content with replacement_text and enforce:
      - Font: Arial
      - Size: 12 pt
    """
    tf = cell.text_frame
    # Clear existing content (paragraphs/runs)
    tf.clear()
    # Create one paragraph with one run
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = replacement_text

    # Apply formatting to the run (you can apply to all runs if you add more)
    font = run.font
    font.name = "Arial"
    font.size = Pt(12)

def replace_text_in_table_cells_on_slide(
    slide,
    target_text: str,
    replacement_text: str,
    *,
    case_sensitive: bool = False,
    exact_match: bool = False,
) -> int:
    """
    Search all table cells on a slide for target_text and replace the entire cell text when found.
    Enforces font name 'Arial' and size 12pt on replaced content.

    Returns:
        int: number of cells replaced.
    """
    replacements = 0
    for shape in slide.shapes:
        for cell in _iter_table_cells_in_shape(shape):
            cell_text = cell.text or ""
            if _text_matches(cell_text, target_text, case_sensitive=case_sensitive, exact_match=exact_match):
                _replace_cell_text_and_format(cell, replacement_text)
                replacements += 1
            elif not exact_match and _text_matches(cell_text, target_text, case_sensitive=case_sensitive, exact_match=False):
                # If substring match is allowed, we still replace the full cell text
                _replace_cell_text_and_format(cell, replacement_text)
                replacements += 1
    return replacements

BLUE_00338D = RGBColor(0x00, 0x33, 0x8D)

def _iter_text_frames_in_shape(shape):
    """
    Yield all text_frame objects reachable from this shape, including:
    - Plain text shapes
    - Placeholders
    - Shapes inside group shapes (recursively)
    - Table cells
    """
    # Group shapes: recurse
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for shp in shape.shapes:
            yield from _iter_text_frames_in_shape(shp)
        return

    # Regular shapes with text
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        yield shape.text_frame

    # Tables
    if getattr(shape, "has_table", False) and shape.has_table:
        tbl = shape.table
        for r in tbl.rows:
            for c in r.cells:
                # Text frame always exists on a cell
                yield c.text_frame

def _text_matches(haystack: str, needle: str, *, case_sensitive: bool, exact_match: bool) -> bool:
    if not case_sensitive:
        haystack, needle = haystack.lower(), needle.lower()
    return (haystack == needle) if exact_match else (needle in haystack)

def _apply_uniform_formatting_to_paragraph(paragraph, *, text: str):
    """
    Replace the paragraph content with 'text' and apply:
    - Arial, 24pt, bold, #00338D
    """
    # Clear all runs and set new text as a single run
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text

    font = run.font
    font.name = "Arial"                # Will often show as "Arial (Headings)" depending on theme
    font.bold = True
    font.size = Pt(24)
    font.color.rgb = BLUE_00338D

def replace_text_on_slide(
    slide,
    target_text: str,
    replacement_text: str,
    *,
    case_sensitive: bool = False,
    exact_match: bool = False,
    replace_entire_shape_when_found: bool = True,
) -> int:
    """
    Search a slide for shapes containing 'target_text'. When found, replace with 'replacement_text'
    and apply formatting (Arial, 24pt, bold, #00338D).

    Parameters:
        slide: pptx.slide.Slide
        target_text: str
        replacement_text: str
        case_sensitive: bool = False
        exact_match: bool = False
            - exact_match=True => shape/paragraph text must exactly equal target_text
            - exact_match=False => substring match is sufficient
        replace_entire_shape_when_found: bool = True
            - True: Replace the entire text of the shape (all paragraphs) with the replacement text.
            - False: Replace per-paragraph where a match is found (paragraph becomes replacement_text).
                     (Still results in uniform formatting in matched paragraphs.)

    Returns:
        int: count of replacements performed (shapes or paragraphs, depending on the mode).
    """
    replacements = 0

    for shape in slide.shapes:
        for tf in _iter_text_frames_in_shape(shape):
            # Gather full text of the shape (for 'entire shape' replacement mode)
            full_text = "\n".join(p.text for p in tf.paragraphs)

            if replace_entire_shape_when_found:
                if _text_matches(full_text, target_text, case_sensitive=case_sensitive, exact_match=exact_match) \
                   or (not exact_match and _text_matches(full_text, target_text, case_sensitive=case_sensitive, exact_match=False)):
                    # Clear entire text frame and replace with one formatted paragraph
                    tf.clear()
                    p = tf.paragraphs[0]
                    _apply_uniform_formatting_to_paragraph(p, text=replacement_text)
                    replacements += 1
            else:
                # Replace per paragraph if it matches or contains the target text
                for p in tf.paragraphs:
                    if _text_matches(p.text, target_text, case_sensitive=case_sensitive, exact_match=exact_match) \
                       or (not exact_match and _text_matches(p.text, target_text, case_sensitive=case_sensitive, exact_match=False)):
                        _apply_uniform_formatting_to_paragraph(p, text=replacement_text)
                        replacements += 1

    return replacements

def duplicate_slide(ppt, slide):
    """
    Duplicate `slide` and append the duplicate to the end of the presentation.
    Returns the new slide object. (Slide positioning is handled by caller.)
    """
    new_slide = ppt.slides.add_slide(slide.slide_layout)
    # remove auto placeholders on new slide
    for shape in list(new_slide.shapes):
        if getattr(shape, "is_placeholder", False):
            sp = shape.element
            sp.getparent().remove(sp)
    # copy shapes (deep XML copy)
    for shape in slide.shapes:
        new_shape = parse_xml(shape.element.xml)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    return new_slide

def update_page_no(ppt, slide_idx, reference_text, replace_text):
    # Find the index of the slide containing the reference text
    part1_idx = None
    for idx, slide in enumerate(ppt.slides):
        if idx == 1:
            continue
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text == reference_text:
                part1_idx = idx
                break
        if part1_idx is not None:
            break

    if part1_idx is None:
        return  # Reference text not found

    # Update the target slide's page number
    target_slide = ppt.slides[slide_idx]
    for shape in target_slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
           if shape.text == replace_text:
                shape.text = str(part1_idx)
                # Format the cell text
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.RIGHT
                    for run in paragraph.runs:
                        run.font.size = Pt(14)
                        run.font.color.rgb = RGBColor(255, 255, 255)   # White

def _to_rgb_color(color):  # Aryan Addition
    """
    Convert a '#RRGGBB' hex string or (R, G, B) tuple into an RGBColor.
    """
    if isinstance(color, str):
        s = color.strip().lstrip('#')
        if len(s) != 6 or any(c not in "0123456789aAbBcCdDeEfF" for c in s):
            raise ValueError(f"Invalid hex color: {color!r}")
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
        return RGBColor(r, g, b)
    elif isinstance(color, (tuple, list)) and len(color) == 3:
        r, g, b = color
        if not all(isinstance(x, int) and 0 <= x <= 255 for x in (r, g, b)):
            raise ValueError(f"Invalid RGB tuple: {color!r}")
        return RGBColor(r, g, b)
    else:
        raise TypeError("color must be '#RRGGBB' or (R, G, B)")

def set_table_cell_format(
    cell,
    *,
    horizontal_align: str | None = "center",   # 'left' | 'center' | 'right' | 'justify' | 'distribute'
    vertical_align: str | None = "middle",     # 'top' | 'middle' | 'bottom'
    font_size_pt: float | None = 12,
    font_color = "#000000",                    # '#RRGGBB' or (R, G, B), None to skip
    bold: bool | None = True,                  # <<< set bold text (True/False/None)
    italic: bool | None = None
):
    """
    Apply alignment, font size, font color, and bold to a python-pptx table cell.

    Parameters
    ----------
    cell : pptx.table._Cell
    horizontal_align : str | None
        'left' | 'center' | 'right' | 'justify' | 'distribute' | None
    vertical_align : str | None
        'top' | 'middle' | 'bottom' | None
    font_size_pt : float | None
        Size in points; None leaves unchanged.
    font_color : str | tuple | None
        '#RRGGBB' or (R,G,B); None leaves unchanged.
    bold : bool | None
        True makes text bold, False removes bold, None leaves unchanged.
    italic : bool | None
        True/False/None similar to bold.
    """
    # --- vertical alignment (cell-level) ---
    if vertical_align:
        va_map = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        key = vertical_align.strip().lower()
        if key not in va_map:
            raise ValueError(f"vertical_align must be one of {list(va_map.keys())}")
        cell.vertical_anchor = va_map[key]

    tf = cell.text_frame  # ensure text frame

    # --- horizontal alignment (paragraph-level) ---
    if horizontal_align:
        ha_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
            "distribute": PP_ALIGN.DISTRIBUTE,
        }
        key = horizontal_align.strip().lower()
        if key not in ha_map:
            raise ValueError(f"horizontal_align must be one of {list(ha_map.keys())}")
        for para in tf.paragraphs:
            para.alignment = ha_map[key]

    # --- font attributes (run-level) ---
    color_rgb = _to_rgb_color(font_color) if font_color is not None else None
    size_pt = Pt(font_size_pt) if font_size_pt is not None else None

    for para in tf.paragraphs:
        # If paragraph has no runs, create one to apply formatting
        if not para.runs:
            para.add_run()
        for run in para.runs:
            font = run.font
            if size_pt is not None:
                font.size = size_pt
            if color_rgb is not None:
                font.color.rgb = color_rgb
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic

def set_table_cell_fill(cell, color, *, set_text_color=None):  #Aryan Addition
    """
    Set a python-pptx table cell's background fill to a solid color.

    Parameters
    ----------
    cell : pptx.table._Cell
        A table cell object (e.g., table.cell(row, col)).
    color : str | tuple[int, int, int]
        Fill color as '#RRGGBB' or (R, G, B).
    set_text_color : None | str | tuple[int, int, int]
        If provided, also set text (font) color in the cell.

    Notes
    -----
    - This uses a SOLID fill. (python-pptx doesn't support gradient fills.)
    - If the cell previously had a theme/gradient/pattern fill, we clear and
      apply a solid fill.
    """
    rgb = _to_rgb_color(color)

    # Ensure fill exists and is set to solid
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = rgb

    # Optionally set text color for all runs in the cell
    if set_text_color is not None:
        txt_rgb = _to_rgb_color(set_text_color)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = txt_rgb

def _fill_flags_and_rr(new_slide, i, raw_df):
    
    # Checkbox (flag) & Risk Rating helpers (initialized once)
    def is_yes(val):
        return (val == '✔') or (str(val).strip().lower() == 'yes')

    flag_map = {
        '*D':  ('D',  8),
        '*O':  ('O',  8),
        '*IT': ('IT', 8),  
        '*FR': ('FR', 8),
        '*OE': ('OE', 8),
        '*C':  ('C',  8),
    }

    rr_map = {
        'H': ('High',   '#FF0000', '#FFFFFF'),
        'M': ('Medium', '#FFC000', '#000000'),
        'L': ('Low',    '#92D050', '#000000'),
    }

    """Apply checkbox and RR logic on the 2nd table (if present)."""
    tables = [sh.table for sh in new_slide.shapes if getattr(sh, "has_table", False)]
    if len(tables) <= 1:
        return
    tbl2 = tables[1]

    for row in tbl2.rows:
        for cell in row.cells:
            txt = cell.text.strip()

            # Flags (*D, *O, *IT, *FR, *OE, *C)
            if txt in flag_map:
                col, font_size = flag_map[txt]
                val = raw_df["Observation listing"].loc[i, col]
                cell.text = '✔' if is_yes(val) else ' '
                set_table_cell_format(
                    cell,
                    horizontal_align="center",
                    vertical_align="middle",
                    font_size_pt=font_size,
                    font_color="#000000",
                    bold=True
                )
                continue

            # Risk Rating (RR)
            if txt == 'RR':
                rr_val = str(raw_df["Observation listing"].loc[i, "RR"]).strip().upper()
                if rr_val in rr_map:
                    label, fill_hex, font_hex = rr_map[rr_val]
                    cell.text = label
                    set_table_cell_fill(cell, fill_hex)
                    set_table_cell_format(
                        cell,
                        horizontal_align="center",
                        vertical_align="middle",
                        font_size_pt=10,
                        font_color=font_hex,
                        bold=True
                    )
                # else leave as is


def delete_slides_by_index(prs, indices):
    """
    Delete slides from a Presentation by 0-based indices.

    Works by removing entries from the slideIdLst and dropping the related part.
    Sorts indices descending to avoid index shifting during removal.
    """
    # Ensure we have a list of slides and sort indices descending
    slides = list(prs.slides)
    max_idx = len(slides) - 1
    to_delete = sorted(set(i for i in indices if 0 <= i <= max_idx), reverse=True)

    if not to_delete:
        return 0

    # XML list of <p:sldId> elements
    sldIdLst = prs.slides._sldIdLst
    sldId_elems = list(sldIdLst)  # iterable of sldId elements

    count = 0
    for idx in to_delete:
        slide = slides[idx]
        # Relationship id of the slide part (e.g., rId256)
        rId = sldId_elems[idx].rId
        # Drop the <p:sldId> element
        sldIdLst.remove(sldId_elems[idx])
        # Drop the relationship to the slide part
        prs.part.drop_rel(rId)
        count += 1

    return count


def fill_detailed_observation_templates(og_df, ppt, template_slide1, template_slide2, slide_idx):
    
    # Running the loop through the df
    raw_df = og_df["Observation listing"]
    for i, row in raw_df.iterrows():
        print(i, row['RR'])

        # Duplicate and append template slide 1
        new_slide1 = duplicate_slide(ppt, template_slide1)
        # Move appended slide to the correct position (right after current template + i)
        sldIdLst = ppt.slides._sldIdLst
        new_id = sldIdLst[-1]   # the most recently appended slide element
        # Re-resolve the template index each time (deck is changing)
        template_index = list(ppt.slides).index(template_slide2)
        insert_pos = slide_idx + i
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_pos, new_id)
        new_slide1 = ppt.slides[insert_pos]  # updated reference

        # Duplicate and append template slide 1
        new_slide2 = duplicate_slide(ppt, template_slide2)
        # Move appended slide to the correct position (right after current template + i)
        sldIdLst = ppt.slides._sldIdLst
        new_id = sldIdLst[-1]   # the most recently appended slide element
        # Re-resolve the template index each time (deck is changing)
        insert_pos = slide_idx + 1 + i
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_pos, new_id)
        new_slide2 = ppt.slides[insert_pos]  # updated reference

        # Working on template slide 1

        # Changing the heading:
        slide1_heading = str(row["Issue heading"]) + " (1/2)"
        replace_text_on_slide(
            new_slide1, "{{Issue heading}}", 
            slide1_heading, 
            case_sensitive=True,
            exact_match=True,
            replace_entire_shape_when_found=True
            )
        
        # Working on table in template
        replace_text_in_table_cells_on_slide(new_slide1, "{{Observations}}", str(row["Observations"]), case_sensitive=True, exact_match= True)
        replace_text_in_table_cells_on_slide(new_slide1, "{{Root cause}}", str(row["Root cause"]), case_sensitive=True, exact_match= True)

        # Working on the Risk Tagging, Root Cause, Risk Type table
        _fill_flags_and_rr(new_slide1, i, og_df)

        # Working on template slide 2

        # Changing the heading:
        slide1_heading = str(row["Issue heading"]) + " (2/2)"
        replace_text_on_slide(
            new_slide2, "{{Issue heading}}", 
            slide1_heading, 
            case_sensitive=True,
            exact_match=True,
            replace_entire_shape_when_found=True
            )
        
        # Working on table in template
        replace_text_in_table_cells_on_slide(new_slide2, "{{Potential risk}}", str(row["Potential risk"]), case_sensitive=True, exact_match= True)
        replace_text_in_table_cells_on_slide(new_slide2, "{{Recommendation}}", str(row["Recommendation"]), case_sensitive=True, exact_match= True)
        replace_text_in_table_cells_on_slide(new_slide2, "{{Management Action Plan}}", str(row["Management Action Plan"]), case_sensitive=True, exact_match= True)

        _fill_flags_and_rr(new_slide2, i, og_df)

        slide_idx += 1
    
    delete_slides_by_index(ppt, [list(ppt.slides).index(template_slide1),list(ppt.slides).index(template_slide2)])

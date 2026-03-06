import re
from pathlib import Path
from pptx import Presentation
import io
import re
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from .Cell_text_Formatting import _parse_color , format_text_in_table_cell
from .Dataframe_mapping import make_ordered_hashmap , map_columns , _find_table_for_mapping , remove_whitespace_cells
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

def _to_rgb_color(color):  # Aryan Addition
    """
    Convert a '#RRGGBB' hex string or (R, G, B) tuple into an RGBColor.
    """
    if isinstance(color, str):
        s = color.strip().lstrip('#')
        if len(s) != 6 or any(c not in "0123456789aAbBcCdDeEfF" for c in s):
            raise ValueError(f"Invalid hex color: {color!r}")
        r = int(s[0:2], 16)
        g = int(s[2:4], 16)
        b = int(s[4:6], 16)
        return RGBColor(r, g, b)
    elif isinstance(color, (tuple, list)) and len(color) == 3:
        r, g, b = color
        if not all(isinstance(x, int) and 0 <= x <= 255 for x in (r, g, b)):
            raise ValueError(f"Invalid RGB tuple: {color!r}")
        return RGBColor(r, g, b)
    else:
        raise TypeError("color must be '#RRGGBB' or (R, G, B)")

def set_table_cell_format(
    cell,
    *,
    horizontal_align: str | None = "center",   # 'left' | 'center' | 'right' | 'justify' | 'distribute'
    vertical_align: str | None = "middle",     # 'top' | 'middle' | 'bottom'
    font_size_pt: float | None = 12,
    font_color = "#000000",                    # '#RRGGBB' or (R, G, B), None to skip
    bold: bool | None = True,                  # <<< set bold text (True/False/None)
    italic: bool | None = None
):
    """
    Apply alignment, font size, font color, and bold to a python-pptx table cell.

    Parameters
    ----------
    cell : pptx.table._Cell
    horizontal_align : str | None
        'left' | 'center' | 'right' | 'justify' | 'distribute' | None
    vertical_align : str | None
        'top' | 'middle' | 'bottom' | None
    font_size_pt : float | None
        Size in points; None leaves unchanged.
    font_color : str | tuple | None
        '#RRGGBB' or (R,G,B); None leaves unchanged.
    bold : bool | None
        True makes text bold, False removes bold, None leaves unchanged.
    italic : bool | None
        True/False/None similar to bold.
    """
    # --- vertical alignment (cell-level) ---
    if vertical_align:
        va_map = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        key = vertical_align.strip().lower()
        if key not in va_map:
            raise ValueError(f"vertical_align must be one of {list(va_map.keys())}")
        cell.vertical_anchor = va_map[key]

    tf = cell.text_frame  # ensure text frame

    # --- horizontal alignment (paragraph-level) ---
    if horizontal_align:
        ha_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
            "distribute": PP_ALIGN.DISTRIBUTE,
        }
        key = horizontal_align.strip().lower()
        if key not in ha_map:
            raise ValueError(f"horizontal_align must be one of {list(ha_map.keys())}")
        for para in tf.paragraphs:
            para.alignment = ha_map[key]

    # --- font attributes (run-level) ---
    color_rgb = _to_rgb_color(font_color) if font_color is not None else None
    size_pt = Pt(font_size_pt) if font_size_pt is not None else None

    for para in tf.paragraphs:
        # If paragraph has no runs, create one to apply formatting
        if not para.runs:
            para.add_run()
        for run in para.runs:
            font = run.font
            if size_pt is not None:
                font.size = size_pt
            if color_rgb is not None:
                font.color.rgb = color_rgb
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic

def set_table_cell_fill(cell, color, *, set_text_color=None):  #Aryan Addition
    """
    Set a python-pptx table cell's background fill to a solid color.

    Parameters
    ----------
    cell : pptx.table._Cell
        A table cell object (e.g., table.cell(row, col)).
    color : str | tuple[int, int, int]
        Fill color as '#RRGGBB' or (R, G, B).
    set_text_color : None | str | tuple[int, int, int]
        If provided, also set text (font) color in the cell.

    Notes
    -----
    - This uses a SOLID fill. (python-pptx doesn't support gradient fills.)
    - If the cell previously had a theme/gradient/pattern fill, we clear and
      apply a solid fill.
    """
    rgb = _to_rgb_color(color)

    # Ensure fill exists and is set to solid
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = rgb

    # Optionally set text color for all runs in the cell
    if set_text_color is not None:
        txt_rgb = _to_rgb_color(set_text_color)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = txt_rgb

def duplicate_slide(ppt, slide):
    """
    Duplicate `slide` and append the duplicate to the end of the presentation.
    Returns the new slide object. (Slide positioning is handled by caller.)
    """
    new_slide = ppt.slides.add_slide(slide.slide_layout)
    # remove auto placeholders on new slide
    for shape in list(new_slide.shapes):
        if getattr(shape, "is_placeholder", False):
            sp = shape.element
            sp.getparent().remove(sp)
    # copy shapes (deep XML copy)
    for shape in slide.shapes:
        new_shape = parse_xml(shape.element.xml)
        new_slide.shapes._spTree.insert_element_before(new_shape, 'p:extLst')
    return new_slide

def split_text_by_max_lines(text, max_lines):
    """
    Return (first_part, remaining) where first_part has at most max_lines lines.
    """
    if text is None:
        return "", ""
    lines = str(text).splitlines()
    if len(lines) <= max_lines:
        return "\n".join(lines), ""
    return "\n".join(lines[:max_lines]), "\n".join(lines[max_lines:])

def handle_cell_overflow(ppt, template_slide, overflow_map,
                         max_lines_per_cell=14,
                         title_shape_idx=0,
                         title_text=None,
                         title_font_pt=14,
                         title_bold=True,
                         cell_font_pt=10,
                         cell_bold=False):
    """
    Process overflow_map: dict[(r, c)] -> remaining_text.
    For each (r,c) create one or more duplicates of template_slide and write
    the remaining text into the same (r,c) cell on subsequent duplicates,
    splitting further if the remaining text still exceeds max_lines_per_cell.

    - ppt: Presentation object
    - template_slide: slide object used as template for duplication
    - overflow_map: dict mapping (row_index, col_index) -> remaining_text (string)
    - other args: formatting and title options for the duplicated slides
    """
    if not overflow_map:
        return

    # helper to duplicate and insert immediately after a given slide
    def _duplicate_and_insert_after(ppt, slide, insert_after_index):
        new_slide = duplicate_slide(ppt, slide)
        sldIdLst = ppt.slides._sldIdLst
        new_id = sldIdLst[-1]
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_after_index + 1, new_id)
        return ppt.slides[insert_after_index + 1]

    # find template index once
    template_index = list(ppt.slides).index(template_slide)

    # process entries in ascending order of keys for deterministic behavior
    # each entry may produce multiple duplicates if remaining text is long
    for (r, c), remaining_text in sorted(overflow_map.items()):
        rem = remaining_text
        # keep duplicating until rem is empty
        while rem:
            # duplicate and insert after the last inserted slide for this template
            # compute current insert position: template_index + 1 + n_inserted
            # count how many slides already follow the template that were created by this function:
            # we approximate by counting slides between template_index and end that have the same layout as template
            # (simpler: always insert immediately after the last slide that was inserted for this template)
            # find last slide index that was inserted after template and has same layout (fallback to template_index)
            last_insert_pos = template_index
            # scan forward to find last slide that was inserted after template (we assume duplicates are contiguous)
            for idx in range(template_index + 1, len(ppt.slides)):
                last_insert_pos = idx

            new_slide = _duplicate_and_insert_after(ppt, template_slide, last_insert_pos)
            # optionally set title text
            if title_text is not None and getattr(new_slide.shapes[title_shape_idx], "text_frame", None) is not None:
                format_text_in_shape(new_slide.shapes[title_shape_idx], title_text, font_pt=title_font_pt, bold=title_bold, align=PP_ALIGN.LEFT)

            # find first table on new_slide
            tbl_sh = None
            for sh in new_slide.shapes:
                if getattr(sh, "has_table", False):
                    tbl_sh = sh
                    break
            if tbl_sh is None:
                # nothing to write into; break to avoid infinite loop
                break
            tbl = tbl_sh.table
            # guard coordinates
            if not (0 <= r < len(tbl.rows) and 0 <= c < len(tbl.columns)):
                break

            # split rem into first_part and leftover
            first_part, rem = split_text_by_max_lines(rem, max_lines_per_cell)
            # write first_part into the target cell on the new slide
            try:
                format_text_in_table_cell(tbl.rows[r].cells[c], first_part, font_pt=cell_font_pt, bold=cell_bold, align=PP_ALIGN.LEFT)
            except Exception:
                # fallback: direct assignment if formatting helper fails
                try:
                    tbl.rows[r].cells[c].text = first_part
                except Exception:
                    pass
            # loop continues if rem still non-empty (another duplicate will be created)

def format_text_in_shape(shape, text, font_pt=12, bold=False, align=PP_ALIGN.LEFT):
    """
    Set `text` into a shape that has a text_frame, applying simple formatting.
    - shape: a pptx shape with .text_frame
    - text: string to insert (keeps newlines)
    - font_pt: font size in points
    - bold: boolean
    - align: PP_ALIGN value
    """
    if getattr(shape, "text_frame", None) is None:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = "" if text is None else str(text)
    run.font.size = Pt(font_pt)
    run.font.bold = bold


def check_overflow_text(text, max_lines=14):
    lines = str(text).splitlines()
    if len(lines) > max_lines:
        return True, "\n".join(lines[:max_lines]), "\n".join(lines[max_lines:])
    return False, "\n".join(lines), ""
def remove_slide(pres, slide):
    """
    Remove `slide` from Presentation `pres`.
    Usage: remove_slide(ppt, template_slide)
    """
    slides = list(pres.slides)
    try:
        idx = slides.index(slide)
    except ValueError:
        # slide not found in presentation
        return
    sldIdLst = pres.slides._sldIdLst
    try:
        sldId = sldIdLst[idx]
        sldIdLst.remove(sldId)
    except Exception:
        # best-effort removal; ignore failures
        pass

def format_text_in_title(shape, text, font_pt, font_name, font_color, bold=False, align=None):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    # normalize font size to points (handle numeric pts, Pt objects, or EMU values)
    if isinstance(font_pt, (int, float)):
        val = float(font_pt)
        pts = val / 12700.0 if val > 1000 else val
    elif hasattr(font_pt, "pt"):
        pts = float(font_pt.pt)
    else:
        pts = float(font_pt)
    font_len = Pt(pts)

    # normalize color
    if isinstance(font_color, tuple) and len(font_color) == 3:
        font_color = RGBColor(*font_color)

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)

    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align if align is not None else PP_ALIGN.LEFT
    p.line_spacing = 1.2

    run = p.runs[0] if p.runs else p.add_run()
    run.text = p.text
    run.font.name = font_name
    run.font.size = font_len
    run.font.bold = bold
    run.font.color.rgb = font_color

#for inserting cell value in templates
def fill_template_with_header_map_and_title(
    raw_df,
    ppt,
    template_slide,
    header_cell_map,
    header_map,
    title_shape_idx=0,
    title_template="{{Issue heading}}",
    title_font_pt=14,
    title_bold=True,
    cell_font_pt=10,
    cell_bold=False ,
    end_char: str = None
):
    """
    Duplicate the template slide once per item in header_map and:
      - set shape[title_shape_idx].text using title_template (replace tokens)
      - write the i-th value of each header key into the table cell specified by header_cell_map[key] = (row, col)
      - apply formatting via format_text_in_shape and format_text_in_table_cell

    Inserts duplicates in ascending order (index 0 first).
    """
    if header_cell_map is None or header_map is None:
        raise ValueError("header_cell_map and header_map must be provided")

    # determine number of slides to create (max length among relevant keys)
    keys_to_consider = set(header_cell_map.keys()) | {"Issue heading", "Sl No.", "Sl No"}
    lengths = []
    for k in keys_to_consider:
        v = header_map.get(k, [])
        try:
            lengths.append(len(v))
        except Exception:
            lengths.append(1 if v not in (None, "") else 0)
    total = max(lengths) if lengths else 0
    if total == 0:
        return

    template_slide = template_slide
    # original index of template in the slides list (used to insert duplicates in order)
    template_index = list(ppt.slides).index(template_slide)

    def _find_first_table(slide):
        for sh in slide.shapes:
            if getattr(sh, "has_table", False):
                return sh
        return None

    for i in range(total):
        # 1) duplicate and append to end
        new_slide = duplicate_slide(ppt, template_slide)

        # 2) move the new slide id element to the correct position so order is preserved
        sldIdLst = ppt.slides._sldIdLst
        new_id = sldIdLst[-1]  # element for the newly appended slide
        # compute desired insertion index: right after template + number of already inserted slides (i)
        insert_pos = template_index + 1 + i
        # remove and insert at desired position
        sldIdLst.remove(new_id)
        sldIdLst.insert(insert_pos, new_id)
        # now obtain the slide object at that position
        new_slide = ppt.slides[insert_pos]  # updated reference to the moved slide

        # 3) set title/heading shape text using template tokens
        try:
            shape0 = new_slide.shapes[title_shape_idx]
        except Exception:
            shape0 = None

        issue_vals = header_map.get("Issue heading", [])
        slno_vals = header_map.get("Sl No.", header_map.get("Sl No", []))

        # pick i-th element if list-like, else use scalar for i==0
        def _pick(vals, idx):
            if vals is None:
                return ""
            if isinstance(vals, (list, tuple)):
                return vals[idx] if idx < len(vals) else ""
            return vals if idx == 0 else ""

        issue_text = _pick(issue_vals, i)
        slno_text = _pick(slno_vals, i)

        title_str = title_template.replace("{{Issue heading}}", str(issue_text).strip()).replace("{{Sl No.}}", str(slno_text).strip())
        title_font_pt = Pt(24)
        title_font_name = "Arial"
        title_font_color = RGBColor(0, 51, 141)
        title_bold = True

        if shape0 is not None and getattr(shape0, "text_frame", None) is not None:
            format_text_in_title(
                shape0,
                title_str,
                font_pt=title_font_pt,
                font_name=title_font_name,
                font_color=title_font_color,
                bold=title_bold,
                align=PP_ALIGN.LEFT
            )

        # if shape0 is not None and getattr(shape0, "text_frame", None) is not None:
        #     format_text_in_shape(shape0, title_str, font_pt=title_font_pt, bold=title_bold, align=PP_ALIGN.LEFT)

        # Fill the Risk Rating column in detailed observations

        tables = []
        for shape in new_slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                tables.append(shape.table)
        
        if len(tables)>1: 
            tbl2 = tables[1]

            for row in tbl2.rows:
                for cell in row.cells:
                    if cell.text == '*D':
                        if raw_df["Observation listing"].loc[i,"D"] == '✔' or str(raw_df["Observation listing"].loc[i,"D"]).lower == 'yes':
                            cell.text = '✔'
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                        
                        else:
                            cell.text = " "
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                    
                    if cell.text == '*O':
                        if raw_df["Observation listing"].loc[i,"O"] == '✔' or str(raw_df["Observation listing"].loc[i,"O"]).lower == 'yes':
                            cell.text = '✔'
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                            
                        else:
                            cell.text = " "
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)

                    if cell.text == '*IT':
                        if raw_df["Observation listing"].loc[i,"IT"] == '✔' or str(raw_df["Observation listing"].loc[i,"IT"]).lower == 'yes':
                            cell.text = '✔'
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=6, font_color="#000000", bold=True)
                            
                        else:
                            cell.text = " "
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)

                    if cell.text == '*FR':
                        if raw_df["Observation listing"].loc[i,"FR"] == '✔' or str(raw_df["Observation listing"].loc[i,"FR"]).lower == 'yes':
                            cell.text = '✔'
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                            
                        else:
                            cell.text = " "
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)

                    if cell.text == '*OE':
                        if raw_df["Observation listing"].loc[i,"OE"] == '✔' or str(raw_df["Observation listing"].loc[i,"OE"]).lower == 'yes':
                            cell.text = '✔'
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                            
                        else:
                            cell.text = " "
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                    
                    if cell.text == '*C':
                        if raw_df["Observation listing"].loc[i,"C"] == '✔' or str(raw_df["Observation listing"].loc[i,"C"]).lower == 'yes':
                            cell.text = '✔'
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)
                            
                        else:
                            cell.text = " "
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=8, font_color="#000000", bold=True)

                    if cell.text == 'RR':
                        if raw_df["Observation listing"].loc[i,"RR"] == 'H':
                            cell.text = 'High'
                            set_table_cell_fill(cell, '#FF0000')
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=10, font_color="#FFFFFF", bold=True)
                        
                        elif raw_df["Observation listing"].loc[i,"RR"] == 'M':
                            cell.text = 'Medium'
                            set_table_cell_fill(cell, '#FFC000')
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=10, font_color="#000000", bold=True)
                        
                        elif raw_df["Observation listing"].loc[i,"RR"] == 'L':
                            cell.text = 'Low'
                            set_table_cell_fill(cell, '#92D050')
                            set_table_cell_format(cell, horizontal_align="center", vertical_align="middle", font_size_pt=10, font_color="#000000", bold=True)
                        
                        else:
                            continue


        # 4) find the first table and fill mapped cells with formatting
        #table_shape = _find_first_table(new_slide)
        table_shape = _find_table_for_mapping(new_slide, header_cell_map)

        if table_shape is None:
            continue
        tbl = table_shape.table
        # inside your loop where you write each cell (replace the simple write with this)
        max_lines_per_cell = 14  # or pass as parameter to your function

        for key, (r, c) in header_cell_map.items():
            vals = header_map.get(key, [])
            text_val = _pick(vals, i)

            # guard against out-of-range coordinates
            if not (0 <= r < len(tbl.rows) and 0 <= c < len(tbl.columns)):
                continue

            # split into allowed and overflow
            allowed, remaining = split_text_by_max_lines(text_val, max_lines_per_cell)
            # write allowed part into current duplicate
            format_text_in_table_cell(tbl.rows[r].cells[c], allowed, font_pt=cell_font_pt, bold=cell_bold, align=PP_ALIGN.LEFT)

            # if there is remaining text, collect into overflow_map to be processed after finishing this slide
            if remaining:
                # use a local overflow_map per slide (or global aggregator)
                local_overflow = {(r, c): remaining}
                # call handler to create duplicates and write remaining text
                # pass the current new_slide (the duplicate you just created) as template for further duplicates
                handle_cell_overflow(ppt, new_slide, local_overflow,
                                    max_lines_per_cell=max_lines_per_cell,
                                    title_shape_idx=title_shape_idx,
                                    title_text=title_str,  # keep same title on overflow slides
                                    title_font_pt=title_font_pt,
                                    title_bold=title_bold,
                                    cell_font_pt=cell_font_pt,
                                    cell_bold=cell_bold)
                

    if end_char:
        ending_slide_index = template_index + total
        ending_slide_number = ending_slide_index + 1
        return ending_slide_number

def fill_templates_with_header_map_and_title_v2(
    raw_df,
    ppt,
    template_slides,               # ← can be a single slide or a list/tuple of slides
    header_cell_map,
    header_map,
    title_shape_idx=0,
    title_template="{{Issue heading}}",
    title_font_pt=14,
    title_bold=True,
    cell_font_pt=10,
    cell_bold=False,
    end_char: str = None,
):
    """
    For each template slide (one or many), duplicates it once per item inferred from header_map and:
      - sets shape[title_shape_idx] using title_template (token replace: {{Issue heading}}, {{Sl No.}})
      - writes values to mapped cells: header_cell_map[key] = (row, col)
      - formats text using format_text_in_title / format_text_in_table_cell
      - applies checkbox flags (*D, *O, *IT, *FR, *OE, *C) and risk rating (RR) logic in the 2nd table (if present)
      - handles overflow via handle_cell_overflow (if remaining text after split)

    Duplicates for each template are inserted immediately after that template, in ascending order.
    """

    if header_cell_map is None or header_map is None:
        raise ValueError("header_cell_map and header_map must be provided")

    # Normalize template_slides to a list
    if not isinstance(template_slides, (list, tuple)):
        template_slides = [template_slides]

    # -------- Derive how many duplicates to create (max length across relevant keys) --------
    keys_to_consider = set(header_cell_map.keys()) | {"Issue heading", "Sl No.", "Sl No"}
    lengths = []
    for k in keys_to_consider:
        v = header_map.get(k, [])
        try:
            lengths.append(len(v))
        except Exception:
            lengths.append(1 if v not in (None, "") else 0)
    total = max(lengths) if lengths else 0
    if total == 0:
        return

    # Pre-fetch header lists for efficiency
    issue_vals = header_map.get("Issue heading", [])
    slno_vals = header_map.get("Sl No.", header_map.get("Sl No", []))

    def _pick(vals, idx):
        """Pick the idx-th value if sequence-like, else scalar for idx==0."""
        if vals is None:
            return ""
        if isinstance(vals, (list, tuple)):
            return vals[idx] if idx < len(vals) else ""
        return vals if idx == 0 else ""

    # Checkbox (flag) & Risk Rating helpers (initialized once)
    def is_yes(val):
        return (val == '✔') or (str(val).strip().lower() == 'yes')

    flag_map = {
        '*D':  ('D',  8),
        '*O':  ('O',  8),
        '*IT': ('IT', 8),  
        '*FR': ('FR', 8),
        '*OE': ('OE', 8),
        '*C':  ('C',  8),
    }

    rr_map = {
        'H': ('High',   '#FF0000', '#FFFFFF'),
        'M': ('Medium', '#FFC000', '#000000'),
        'L': ('Low',    '#92D050', '#000000'),
    }

    def _fill_flags_and_rr(new_slide, i):
        """Apply checkbox and RR logic on the 2nd table (if present)."""
        tables = [sh.table for sh in new_slide.shapes if getattr(sh, "has_table", False)]
        if len(tables) <= 1:
            return
        tbl2 = tables[1]

        for row in tbl2.rows:
            for cell in row.cells:
                txt = cell.text.strip()

                # Flags (*D, *O, *IT, *FR, *OE, *C)
                if txt in flag_map:
                    col, font_size = flag_map[txt]
                    val = raw_df["Observation listing"].loc[i, col]
                    cell.text = '✔' if is_yes(val) else ' '
                    set_table_cell_format(
                        cell,
                        horizontal_align="center",
                        vertical_align="middle",
                        font_size_pt=font_size,
                        font_color="#000000",
                        bold=True
                    )
                    continue

                # Risk Rating (RR)
                if txt == 'RR':
                    rr_val = str(raw_df["Observation listing"].loc[i, "RR"]).strip().upper()
                    if rr_val in rr_map:
                        label, fill_hex, font_hex = rr_map[rr_val]
                        cell.text = label
                        set_table_cell_fill(cell, fill_hex)
                        set_table_cell_format(
                            cell,
                            horizontal_align="center",
                            vertical_align="middle",
                            font_size_pt=10,
                            font_color=font_hex,
                            bold=True
                        )
                    # else leave as is

    # Fallback if _find_table_for_mapping isn't present in your environment
    def _find_first_table(slide):
        for sh in slide.shapes:
            if getattr(sh, "has_table", False):
                return sh
        return None

    # ------------------------- MAIN: per template slide -------------------------
    ending_slide_number = None  # will hold the last resulting "human index" if requested

    for template_slide in template_slides:
        # For each template, insert duplicates *right after* it in ascending order
        for i in range(total):
            # Duplicate and append
            new_slide = duplicate_slide(ppt, template_slide)

            # Move appended slide to the correct position (right after current template + i)
            sldIdLst = ppt.slides._sldIdLst
            new_id = sldIdLst[-1]   # the most recently appended slide element
            # Re-resolve the template index each time (deck is changing)
            template_index = list(ppt.slides).index(template_slide)
            insert_pos = template_index + 1 + i
            sldIdLst.remove(new_id)
            sldIdLst.insert(insert_pos, new_id)
            new_slide = ppt.slides[insert_pos]  # updated reference

            # ---- Title / Heading ----
            issue_text = _pick(issue_vals, i)
            slno_text  = _pick(slno_vals,  i)

            title_str = (
                title_template
                .replace("{{Issue heading}}", str(issue_text).strip())
                .replace("{{Sl No.}}",       str(slno_text).strip())
            )

            # Allow caller-specified font/weight; default left align
            try:
                shape0 = new_slide.shapes[title_shape_idx]
            except Exception:
                shape0 = None

            if shape0 is not None and getattr(shape0, "text_frame", None) is not None:
                # Prefer your richer title formatter if available
                try:
                    format_text_in_title(
                        shape0,
                        title_str,
                        font_pt=title_font_pt,
                        # If you want a default font/color, set here; else keep existing
                        font_name="Arial", 
                        font_color=RGBColor(0, 51, 141),
                        bold=title_bold,
                        align=PP_ALIGN.LEFT
                    )
                except NameError:
                    # Fallback to a generic formatter if format_text_in_title isn't available
                    format_text_in_shape(
                        shape0, title_str,
                        font_pt=title_font_pt,
                        bold=title_bold,
                        align=PP_ALIGN.LEFT
                    )

            # ---- Flags & RR on the 2nd table (if present) ----
            _fill_flags_and_rr(new_slide, i)

            # ---- Fill mapped cells in the target table ----
            try:
                table_shape = _find_table_for_mapping(new_slide, header_cell_map)
            except NameError:
                table_shape = _find_first_table(new_slide)

            if table_shape is None:
                continue

            tbl = table_shape.table
            max_lines_per_cell = 14  # can be parameterized

            for key, (r, c) in header_cell_map.items():
                vals = header_map.get(key, [])
                text_val = _pick(vals, i)

                # Guard against out-of-range (defensive)
                if not (0 <= r < len(tbl.rows) and 0 <= c < len(tbl.columns)):
                    continue

                allowed, remaining = split_text_by_max_lines(text_val, max_lines_per_cell)

                format_text_in_table_cell(
                    tbl.rows[r].cells[c],
                    allowed,
                    font_pt=cell_font_pt,
                    bold=cell_bold,
                    align=PP_ALIGN.LEFT
                )

                if remaining:
                    local_overflow = {(r, c): remaining}
                    handle_cell_overflow(
                        ppt,
                        new_slide,  # use the just-created slide as template for overflow slides
                        local_overflow,
                        max_lines_per_cell=max_lines_per_cell,
                        title_shape_idx=title_shape_idx,
                        title_text=title_str,     # keep same title on overflow slides
                        title_font_pt=title_font_pt,
                        title_bold=title_bold,
                        cell_font_pt=cell_font_pt,
                        cell_bold=cell_bold
                    )

        # For end_char behavior: report the 1-based slide number just after the last insertion for THIS template
        if end_char:
            # After all insertions for this template, compute the “ending” position:
            template_index_now = list(ppt.slides).index(template_slide)
            ending_slide_index = template_index_now + total
            ending_slide_number = ending_slide_index + 1  # human-readable index

    if end_char:
        return ending_slide_number
